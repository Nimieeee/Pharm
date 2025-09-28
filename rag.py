"""
RAG (Retrieval-Augmented Generation) System
Handles document processing and vector search using Langchain
"""

import os
from typing import List, Dict, Any, Optional
import streamlit as st
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import OpenAIEmbeddings
from langchain_community.document_loaders import PyPDFLoader, TextLoader
from langchain.schema import Document
from sentence_transformers import SentenceTransformer
import tempfile
import numpy as np

# Additional imports for enhanced document processing
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# OCR functionality removed

# LangExtract removed for performance

from database import SimpleChatbotDB


class RAGManager:
    """Manages document processing and retrieval for RAG system"""
    
    VERSION = "2.0.0"  # Version to force cache refresh
    
    def __init__(self):
        self.db_manager = SimpleChatbotDB()
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1500,  # Increased from 1000 for larger chunks
            chunk_overlap=300,  # Increased overlap for better continuity
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""]  # Better splitting on paragraphs and sentences
        )
        self.embedding_model = None

        self._initialize_embeddings()
        

    
    def _initialize_embeddings(self):
        """Initialize embedding model - prioritize SentenceTransformers"""
        try:
            # Primary: Use SentenceTransformers with PyTorch fix
            self._initialize_sentence_transformers()
            
            # Only try OpenAI if SentenceTransformers failed and API key is available
            if not self.embedding_model:
                openai_api_key = os.getenv("OPENAI_API_KEY")
                if openai_api_key:
                    try:
                        self.embedding_model = OpenAIEmbeddings(openai_api_key=openai_api_key)
                        st.info("✅ Using OpenAI embeddings as fallback")
                    except Exception as openai_e:
                        st.warning(f"OpenAI embeddings also failed: {str(openai_e)}")
                
        except Exception as e:
            st.error(f"Error initializing embeddings: {str(e)}")
            self.embedding_model = None
    
    def _initialize_sentence_transformers(self):
        """Initialize sentence transformers with PyTorch tensor fix"""
        import torch
        
        # Fix for PyTorch meta tensor issue
        strategies = [
            # Strategy 1: Force CPU and disable meta device
            lambda: self._create_sentence_transformer_safe('all-MiniLM-L6-v2', device='cpu'),
            
            # Strategy 2: With explicit torch settings
            lambda: self._create_sentence_transformer_with_torch_fix('all-MiniLM-L6-v2'),
            
            # Strategy 3: Alternative model path
            lambda: self._create_sentence_transformer_safe('sentence-transformers/all-MiniLM-L6-v2', device='cpu'),
            
            # Strategy 4: Different model entirely
            lambda: self._create_sentence_transformer_safe('all-MiniLM-L12-v2', device='cpu'),
        ]
        
        for i, strategy in enumerate(strategies, 1):
            try:
                self.embedding_model = strategy()
                st.success(f"✅ SentenceTransformer initialized successfully (method {i})")
                return
            except Exception as e:
                error_msg = str(e)
                if i < len(strategies):
                    # Only show brief error for intermediate failures
                    pass
                else:
                    st.error(f"SentenceTransformer initialization failed: {error_msg}")
        
        # If all strategies fail, set to None
        self.embedding_model = None
        st.warning("⚠️ SentenceTransformer initialization failed. Try updating: pip install --upgrade sentence-transformers torch")
    
    def _create_sentence_transformer_safe(self, model_name: str, device: str = 'cpu'):
        """Create SentenceTransformer with safe tensor handling"""
        import torch
        
        # Set torch to avoid meta tensors
        torch.set_default_device(device)
        
        # Create model with explicit device
        model = SentenceTransformer(model_name, device=device)
        
        # Ensure all parameters are on the correct device
        model = model.to(device)
        
        return model
    
    def _create_sentence_transformer_with_torch_fix(self, model_name: str):
        """Create SentenceTransformer with PyTorch tensor fix"""
        import torch
        
        # Temporarily disable meta device usage
        original_device = torch.get_default_device() if hasattr(torch, 'get_default_device') else None
        
        try:
            # Force CPU as default device
            if hasattr(torch, 'set_default_device'):
                torch.set_default_device('cpu')
            
            # Create model
            model = SentenceTransformer(model_name, device='cpu')
            
            # Explicitly move all parameters to CPU using to_empty() if available
            for param in model.parameters():
                if param.is_meta:
                    # Use to_empty() for meta tensors
                    param.data = torch.empty_like(param, device='cpu')
            
            return model
            
        finally:
            # Restore original device setting
            if original_device and hasattr(torch, 'set_default_device'):
                torch.set_default_device(original_device)
        

        

    
    def _initialize_document_intelligence(self):
        """Compatibility method for cached sessions - no longer needed"""
        pass
    
    def process_uploaded_file(self, uploaded_file, conversation_id: str, user_session_id: str, progress_callback=None) -> tuple[bool, int]:
        """
        Process uploaded file and store chunks in database with enhanced feedback
        
        Args:
            uploaded_file: Streamlit uploaded file object
            conversation_id: ID of the conversation to associate chunks with
            user_session_id: ID of the user session for privacy isolation
            progress_callback: Optional callback function for progress updates
            
        Returns:
            Tuple of (success: bool, chunk_count: int)
        """
        try:
            if not uploaded_file:
                return False, 0
            
            if progress_callback:
                progress_callback("Saving temporary file...")
            
            # Save uploaded file temporarily
            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{uploaded_file.name.split('.')[-1]}") as tmp_file:
                tmp_file.write(uploaded_file.getvalue())
                tmp_file_path = tmp_file.name
            
            if progress_callback:
                progress_callback("Loading document...")
            
            # Load document based on file type
            documents = self._load_document(tmp_file_path, uploaded_file.name)
            
            if not documents:
                os.unlink(tmp_file_path)
                return False, 0
            
            if progress_callback:
                progress_callback("Splitting into chunks...")
            
            # Fast document processing without LangExtract
            
            # Split documents into chunks
            chunks = self.text_splitter.split_documents(documents)
            
            if not chunks:
                st.error(f"No chunks created from '{uploaded_file.name}' - document may be empty or unreadable")
                os.unlink(tmp_file_path)
                return False, 0
            
            if progress_callback:
                progress_callback(f"Processing {len(chunks)} chunks...")
            
            # Process and store chunks with progress tracking
            success_count = 0
            for i, chunk in enumerate(chunks):
                # Check if chunk has content
                if not chunk.page_content.strip():
                    continue
                
                if self._process_and_store_chunk(chunk, uploaded_file.name, conversation_id, user_session_id):
                    success_count += 1
            
            # Clean up temporary file
            os.unlink(tmp_file_path)
            
            return success_count > 0, success_count
            
        except Exception as e:
            st.error(f"Error processing file: {str(e)}")
            return False, 0
    
    def _load_document(self, file_path: str, filename: str) -> List[Document]:
        """Load document based on file type with enhanced support including PowerPoint and OCR"""
        try:
            file_extension = filename.lower().split('.')[-1]
            
            if file_extension == 'pdf':
                loader = PyPDFLoader(file_path)
                documents = loader.load()
                if not documents:
                    st.warning(f"PDF '{filename}' loaded but contains no readable content")
                return documents
                
            elif file_extension in ['txt', 'md']:
                loader = TextLoader(file_path, encoding='utf-8')
                documents = loader.load()
                if not documents:
                    st.warning(f"Text file '{filename}' loaded but contains no content")
                return documents
                
            elif file_extension == 'docx':
                return self._process_docx(file_path, filename)
                
            elif file_extension == 'pptx':
                return self._process_pptx(file_path, filename)
                
            elif file_extension in ['jpg', 'jpeg', 'png', 'bmp', 'tiff', 'gif']:
                st.error(f"Image files are not supported: {filename}")
                st.info("💡 Please convert images to PDF or text format")
                return []
                    
            else:
                st.error(f"Unsupported file type: {file_extension}. Supported types: PDF, TXT, MD, DOCX, PPTX")
                return []
                
        except Exception as e:
            st.error(f"Error loading document '{filename}': {str(e)}")
            return []
    
    def _process_docx(self, file_path: str, filename: str) -> List[Document]:
        """Process DOCX files"""
        try:
            import docx
        except ImportError:
            st.error("python-docx package not found. Please install it with: pip install python-docx")
            return []
        
        try:
            doc = docx.Document(file_path)
            
            # Extract text from all paragraphs
            text_content = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text.strip())
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_content.append(cell.text.strip())
            
            if not text_content:
                st.warning(f"DOCX '{filename}' contains no readable text")
                return []
            
            # Create a Document object
            full_text = '\n\n'.join(text_content)
            document = Document(
                page_content=full_text,
                metadata={"source": filename, "file_type": "docx"}
            )
            
            return [document]
            
        except Exception as docx_error:
            st.error(f"❌ Error processing DOCX '{filename}': {str(docx_error)}")
            return []
    
    def _process_pptx(self, file_path: str, filename: str) -> List[Document]:
        """Process PowerPoint files"""
        if not PPTX_AVAILABLE:
            st.error("python-pptx package not found. Please install it with: pip install python-pptx")
            return []
        
        try:
            prs = Presentation(file_path)
            
            text_content = []
            slide_number = 0
            
            for slide in prs.slides:
                slide_number += 1
                slide_text = []
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    slide_content = f"Slide {slide_number}:\n" + '\n'.join(slide_text)
                    text_content.append(slide_content)
            
            if not text_content:
                st.warning(f"PowerPoint '{filename}' contains no readable text")
                return []
            
            # Create a Document object
            full_text = '\n\n'.join(text_content)
            document = Document(
                page_content=full_text,
                metadata={"source": filename, "file_type": "pptx", "slides": slide_number}
            )
            
            return [document]
            
        except Exception as pptx_error:
            st.error(f"❌ Error processing PowerPoint '{filename}': {str(pptx_error)}")
            return []
    

    

    
    def _process_and_store_chunk(self, chunk: Document, filename: str, conversation_id: str, user_session_id: str) -> bool:
        """Process chunk and store in database"""
        try:
            # Generate embedding (required for proper vector search)
            embedding = self._generate_embedding(chunk.page_content)
            if not embedding:
                # Skip this chunk if we can't generate embeddings
                return False
            
            # Prepare metadata
            metadata = {
                "filename": filename,
                "page": chunk.metadata.get("page", 0),
                "source": chunk.metadata.get("source", filename)
            }
            
            # Store in database
            return self.db_manager.store_document_chunk(
                content=chunk.page_content,
                embedding=embedding,
                metadata=metadata,
                conversation_id=conversation_id,
                user_session_id=user_session_id
            )
            
        except Exception as e:
            # Only show error once per session to avoid spam
            if not hasattr(st.session_state, 'chunk_processing_error_shown'):
                st.error(f"Error processing chunk: {str(e)}")
                st.session_state.chunk_processing_error_shown = True
            return False
    
    def _generate_embedding(self, text: str) -> Optional[List[float]]:
        """Generate embedding for text with detailed error handling"""
        try:
            if not self.embedding_model:
                # Don't show error repeatedly - just return None silently
                return None
            
            if not text or not text.strip():
                return None
            
            if hasattr(self.embedding_model, 'embed_query'):
                # OpenAI embeddings
                embedding = self.embedding_model.embed_query(text)
            else:
                # Sentence transformers
                embedding = self.embedding_model.encode(text).tolist()
            
            if not embedding or len(embedding) == 0:
                return None
            
            return embedding
            
        except Exception as e:
            # Only show error once per session to avoid spam
            if not hasattr(st.session_state, 'embedding_error_shown'):
                st.error(f"Error generating embedding: {str(e)}")
                if "api" in str(e).lower():
                    st.info("💡 Check your OpenAI API key or try using sentence-transformers")
                st.session_state.embedding_error_shown = True
            return None
    
    def search_relevant_context(self, query: str, conversation_id: str, user_session_id: str, max_chunks: int = None, include_document_overview: bool = False, unlimited_context: bool = False) -> str:
        """
        Search for relevant context based on query within a conversation and user session
        
        Args:
            query: User query
            conversation_id: ID of the conversation to search within
            user_session_id: ID of the user session for privacy isolation
            max_chunks: Maximum number of chunks to retrieve (None for unlimited)
            include_document_overview: Whether to include document overview/summary
            unlimited_context: Whether to retrieve all available chunks
            
        Returns:
            Concatenated relevant context
        """
        try:
            # Generate query embedding
            query_embedding = self._generate_embedding(query)
            if not query_embedding:
                # If no embeddings available, return all available context
                st.warning("⚠️ Vector search not available, returning all document context")
                return self.get_all_document_context(conversation_id, user_session_id)
            
            # Search for similar chunks - use unlimited if requested
            search_limit = 1000 if unlimited_context or max_chunks is None else (max_chunks or 10)
            similar_chunks = self.db_manager.search_similar_chunks(query_embedding, conversation_id, user_session_id, search_limit, threshold=0.1)
            
            # If no similar chunks found, try getting any chunks as fallback
            if not similar_chunks:
                st.write("Debug - No similar chunks found, trying fallback query...")
                fallback_chunks = self.db_manager.get_random_chunks(conversation_id, max_chunks)
                st.write(f"Debug - Fallback found {len(fallback_chunks)} chunks")
                if fallback_chunks:
                    similar_chunks = fallback_chunks
            
            # For now, always use fallback until vector search is fixed
            if len(similar_chunks) == 0:
                st.write("Debug - Using fallback chunks for context...")
                if unlimited_context:
                    st.write("Debug - Getting ALL chunks for unlimited context...")
                    similar_chunks = self.db_manager.get_all_conversation_chunks(conversation_id, user_session_id)
                else:
                    fallback_limit = max_chunks or 10
                    similar_chunks = self.db_manager.get_random_chunks(conversation_id, user_session_id, fallback_limit)
            
            if not similar_chunks:
                return ""
            
            # Check if user is asking for comprehensive document information
            comprehensive_keywords = [
                "entire document", "whole document", "complete document", "full document",
                "summarize", "summary", "overview", "explain the document", "what is this document about",
                "document content", "all information", "everything in", "comprehensive"
            ]
            
            is_comprehensive_query = any(keyword in query.lower() for keyword in comprehensive_keywords)
            
            if is_comprehensive_query or include_document_overview:
                # For comprehensive queries, get more chunks and add document metadata
                # Comprehensive query detected, retrieving more context
                
                # Get additional chunks for comprehensive coverage
                additional_limit = 1000 if unlimited_context else (max_chunks * 3 if max_chunks else 50)
                additional_chunks = self.db_manager.get_random_chunks(conversation_id, user_session_id, additional_limit)
                
                # Combine and deduplicate chunks
                all_chunks = similar_chunks + additional_chunks
                seen_content = set()
                unique_chunks = []
                
                for chunk in all_chunks:
                    content = chunk.get("content", "")
                    if content and content not in seen_content:
                        seen_content.add(content)
                        unique_chunks.append(chunk)
                
                # Use all unique chunks if unlimited, otherwise limit
                if unlimited_context or max_chunks is None:
                    similar_chunks = unique_chunks  # No limit
                else:
                    similar_chunks = unique_chunks[:max_chunks * 2]
            
            # Concatenate relevant content
            context_parts = []
            document_info = {}
            
            for i, chunk in enumerate(similar_chunks):
                content = chunk.get("content", "")
                similarity = chunk.get("similarity", 0)
                metadata = chunk.get("metadata", {})
                
                # Processing chunk content
                
                if content:
                    context_parts.append(content)
                    
                    # Collect document metadata
                    filename = metadata.get("filename", "Unknown")
                    if filename not in document_info:
                        document_info[filename] = {
                            "chunks": 0,
                            "file_type": metadata.get("file_type", "unknown")
                        }
                    document_info[filename]["chunks"] += 1
            
            # Add document overview if comprehensive query
            if is_comprehensive_query and document_info:
                overview_parts = ["=== DOCUMENT OVERVIEW ==="]
                for filename, info in document_info.items():
                    overview_parts.append(f"Document: {filename} ({info['file_type']}) - {info['chunks']} sections included")
                overview_parts.append("=== DOCUMENT CONTENT ===")
                context_parts = overview_parts + context_parts
            
            # Smart context management - check total length and truncate if needed
            full_context = "\n\n".join(context_parts)
            
            # Mistral models can handle up to ~32k tokens, but let's be conservative
            max_context_chars = 100000  # ~25k tokens approximately
            
            if len(full_context) > max_context_chars:
                # Context too large, truncating for performance
                
                # Prioritize chunks by similarity score
                sorted_chunks = sorted(similar_chunks, key=lambda x: x.get("similarity", 0), reverse=True)
                
                truncated_parts = []
                current_length = 0
                
                # Add overview first if it exists
                if is_comprehensive_query and document_info:
                    for part in overview_parts:
                        if current_length + len(part) < max_context_chars:
                            truncated_parts.append(part)
                            current_length += len(part) + 2  # +2 for \n\n
                
                # Add chunks in order of similarity
                for chunk in sorted_chunks:
                    content = chunk.get("content", "")
                    if content and current_length + len(content) < max_context_chars:
                        truncated_parts.append(content)
                        current_length += len(content) + 2
                    else:
                        break
                
                full_context = "\n\n".join(truncated_parts)
                # Context truncated for optimal performance
            
            # Context is ready for use
            
            return full_context
            
        except Exception as e:
            st.error(f"Error searching context: {str(e)}")
            return ""
    

    
    def get_all_document_context(self, conversation_id: str, user_session_id: str) -> str:
        """
        Get ALL document content for complete context (no chunking limitations)
        
        Args:
            conversation_id: ID of the conversation
            user_session_id: ID of the user session for privacy isolation
            
        Returns:
            Complete document content as context
        """
        try:
            # Get ALL chunks for this conversation and user session ONLY
            all_chunks = self.db_manager.get_all_conversation_chunks(conversation_id, user_session_id)
            
            # Verify conversation isolation - double-check that all chunks belong to this conversation
            verified_chunks = []
            for chunk in all_chunks:
                chunk_conv_id = chunk.get("conversation_id")
                chunk_user_id = chunk.get("user_session_id") 
                
                # Only include chunks that definitely belong to this conversation and user
                if chunk_conv_id == conversation_id and chunk_user_id == user_session_id:
                    verified_chunks.append(chunk)
            
            all_chunks = verified_chunks
            
            if not all_chunks:
                return ""
            
            # Organize chunks by document/filename
            documents = {}
            for chunk in all_chunks:
                metadata = chunk.get("metadata", {})
                filename = metadata.get("filename", "Unknown Document")
                
                if filename not in documents:
                    documents[filename] = {
                        "chunks": [],
                        "file_type": metadata.get("file_type", "unknown")
                    }
                
                documents[filename]["chunks"].append({
                    "content": chunk.get("content", ""),
                    "created_at": chunk.get("created_at", "")
                })
            
            # Build complete context with document organization
            context_parts = []
            
            # Add document overview
            if len(documents) > 1:
                context_parts.append("=== UPLOADED DOCUMENTS OVERVIEW ===")
                for filename, doc_info in documents.items():
                    context_parts.append(f"📄 {filename} ({doc_info['file_type']}) - {len(doc_info['chunks'])} sections")
                context_parts.append("")
            
            # Add complete content for each document
            for filename, doc_info in documents.items():
                if len(documents) > 1:
                    context_parts.append(f"=== DOCUMENT: {filename} ===")
                
                # Sort chunks by creation time to maintain document order
                sorted_chunks = sorted(doc_info["chunks"], key=lambda x: x.get("created_at", ""))
                
                for chunk in sorted_chunks:
                    content = chunk["content"].strip()
                    if content:
                        context_parts.append(content)
                
                if len(documents) > 1:
                    context_parts.append("")  # Separator between documents
            
            # Combine all content
            full_context = "\n\n".join(context_parts)
            
            # Smart truncation if context is too large (keep most important parts)
            max_context_chars = 150000  # ~37k tokens - very generous limit
            
            if len(full_context) > max_context_chars:
                
                # Keep document overview and truncate content proportionally
                if len(documents) > 1:
                    # Find where content starts (after overview)
                    overview_end = full_context.find("=== DOCUMENT:")
                    if overview_end > 0:
                        overview = full_context[:overview_end]
                        content_section = full_context[overview_end:]
                        
                        # Truncate content section
                        available_chars = max_context_chars - len(overview)
                        if available_chars > 0:
                            truncated_content = content_section[:available_chars]
                            full_context = overview + truncated_content
                        else:
                            full_context = overview
                    else:
                        full_context = full_context[:max_context_chars]
                else:
                    full_context = full_context[:max_context_chars]
                
            return full_context
            
        except Exception as e:
            st.error(f"Error getting complete document context: {str(e)}")
            return ""
    
    def get_conversation_document_stats(self, conversation_id: str, user_session_id: str) -> Dict[str, Any]:
        """Get document statistics for a specific conversation and user session"""
        try:
            all_chunks = self.db_manager.get_all_conversation_chunks(conversation_id, user_session_id)
            
            if not all_chunks:
                return {
                    "total_chunks": 0,
                    "unique_documents": 0,
                    "total_size_mb": 0,
                    "documents": []
                }
            
            # Process chunks to get document info
            documents = {}
            total_content_size = 0
            
            for chunk in all_chunks:
                metadata = chunk.get("metadata", {})
                filename = metadata.get("filename", "Unknown")
                content = chunk.get("content", "")
                
                total_content_size += len(content.encode('utf-8'))
                
                if filename not in documents:
                    documents[filename] = {
                        "chunks": 0,
                        "file_type": metadata.get("file_type", "unknown")
                    }
                
                documents[filename]["chunks"] += 1
            
            return {
                "total_chunks": len(all_chunks),
                "unique_documents": len(documents),
                "total_size_mb": round(total_content_size / (1024 * 1024), 2),
                "documents": [{"name": name, **info} for name, info in documents.items()]
            }
            
        except Exception as e:
            st.error(f"❌ Error getting conversation document stats: {str(e)}")
            return {"total_chunks": 0, "unique_documents": 0, "total_size_mb": 0, "documents": []}
    
    def get_document_stats(self) -> Dict[str, Any]:
        """Get comprehensive statistics about stored documents"""
        try:
            chunk_count = self.db_manager.get_chunk_count()
            document_info = self.db_manager.get_document_info()
            
            return {
                "total_chunks": chunk_count,
                "embedding_model": type(self.embedding_model).__name__ if self.embedding_model else "None",
                "unique_documents": document_info.get("unique_documents", 0),
                "total_size_mb": document_info.get("total_size_mb", 0),
                "last_updated": document_info.get("last_updated", "Never")
            }
            
        except Exception as e:
            st.error(f"Error getting document stats: {str(e)}")
            return {
                "total_chunks": 0, 
                "embedding_model": "Error",
                "unique_documents": 0,
                "total_size_mb": 0,
                "last_updated": "Error"
            }
    
    def clear_all_documents(self) -> bool:
        """Clear all stored documents"""
        try:
            return self.db_manager.clear_all_chunks()
        except Exception as e:
            st.error(f"Error clearing documents: {str(e)}")
            return False
    
    def get_processing_status(self) -> Dict[str, Any]:
        """Get current processing status"""
        return getattr(st.session_state, 'processing_status', {
            'active': False,
            'progress': 0,
            'current_file': '',
            'total_files': 0,
            'processed_files': 0
        })
    
    def get_embedding_status(self) -> Dict[str, Any]:
        """Get embedding model status for diagnostics"""
        return {
            'model_available': self.embedding_model is not None,
            'model_type': type(self.embedding_model).__name__ if self.embedding_model else 'None',
            'is_sentence_transformer': 'SentenceTransformer' in type(self.embedding_model).__name__ if self.embedding_model else False,
            'search_method': 'Vector similarity search' if self.embedding_model else 'No embeddings available'
        }