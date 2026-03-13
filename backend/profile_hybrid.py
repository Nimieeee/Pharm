"""
Profile the HYBRID document processing pipeline.
Tests PyMuPDF text extraction + selective VLM for images using Mistral API.
Usage: /var/www/pharmgpt-backend/venv/bin/python profile_hybrid.py <filepath>
"""
import asyncio, time, os, sys, io, base64

async def main():
    filepath = sys.argv[1] if len(sys.argv) > 1 else None
    if not filepath or not os.path.exists(filepath):
        print(f"Usage: python profile_hybrid.py <filepath>"); return

    ext = filepath.rsplit('.', 1)[-1].lower()
    print(f"=== HYBRID PROFILING: {os.path.basename(filepath)} ({ext}) ===\n", flush=True)
    t_total = time.time()

    # Read file
    with open(filepath, "rb") as f:
        content = f.read()
    print(f"[1] File read: {len(content)} bytes", flush=True)

    # Load Mistral API key
    MISTRAL_KEY = os.getenv("MISTRAL_API_KEY", "")
    if not MISTRAL_KEY:
        for ep in ["/var/www/pharmgpt-backend/backend/.env", os.path.join(os.path.dirname(filepath), ".env")]:
            if os.path.exists(ep):
                for line in open(ep):
                    if line.startswith("MISTRAL_API_KEY="):
                        MISTRAL_KEY = line.split("=", 1)[1].strip().strip('"').strip("'")
                        break
            if MISTRAL_KEY: break
    if not MISTRAL_KEY:
        print("ERROR: No MISTRAL_API_KEY found"); return

    if ext == 'pdf':
        await profile_pdf(content, MISTRAL_KEY)
    elif ext == 'pptx':
        await profile_pptx(content, MISTRAL_KEY)
    else:
        print(f"Unsupported: {ext}")

    print(f"\n{'='*50}")
    print(f"TOTAL: {time.time()-t_total:.2f}s", flush=True)


async def profile_pdf(content, api_key):
    import fitz
    from PIL import Image

    t0 = time.time()
    doc = fitz.open(stream=content, filetype="pdf")
    
    text_pages = []
    images_to_analyze = []

    for page_idx in range(len(doc)):
        page = doc[page_idx]
        text = page.get_text("text").strip()
        
        # Extract embedded images
        page_images = page.get_images(full=True)
        for img_idx, img_info in enumerate(page_images):
            try:
                xref = img_info[0]
                img_data = doc.extract_image(xref)
                if img_data and img_data.get("image") and len(img_data["image"]) > 50000:
                    w = img_data.get("width", 0)
                    h = img_data.get("height", 0)
                    if w >= 200 and h >= 200:
                        images_to_analyze.append((page_idx, img_data["image"], f"P{page_idx+1}_Img{img_idx+1}"))
            except:
                pass

        if len(text) >= 100:
            text_pages.append(page_idx)
        else:
            # Scanned page — render for VLM
            try:
                pix = page.get_pixmap(dpi=150)
                images_to_analyze.append((page_idx, pix.tobytes("jpeg"), f"P{page_idx+1}_scan"))
            except:
                pass

    doc.close()
    t_extract = time.time() - t0
    print(f"[2] PyMuPDF extraction: {t_extract:.2f}s", flush=True)
    print(f"    Text pages: {len(text_pages)}/{len(text_pages) + len([i for i in images_to_analyze if 'scan' in i[2]])}", flush=True)
    print(f"    Images to analyze: {len(images_to_analyze)}", flush=True)

    if images_to_analyze:
        await profile_vlm(images_to_analyze, api_key)
    else:
        print(f"[3] VLM: SKIPPED (no images needed)", flush=True)


async def profile_pptx(content, api_key):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    t0 = time.time()
    prs = Presentation(io.BytesIO(content))
    
    slide_count = 0
    images_to_analyze = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_count += 1
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t: texts.append(t)
            if shape.has_table:
                texts.append("[table]")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image.blob
                    if len(img) > 5000:
                        images_to_analyze.append((slide_idx, img, f"S{slide_idx+1}_img"))
                except:
                    pass
        
        total_text = " ".join(texts)
        print(f"    Slide {slide_idx+1}: {len(total_text)} chars text, "
              f"{'has images' if any(i[0]==slide_idx for i in images_to_analyze) else 'text only'}", flush=True)

    t_extract = time.time() - t0
    print(f"[2] python-pptx extraction: {t_extract:.2f}s", flush=True)
    print(f"    Slides: {slide_count}", flush=True)
    print(f"    Images to analyze: {len(images_to_analyze)}", flush=True)

    if images_to_analyze:
        await profile_vlm(images_to_analyze, api_key)
    else:
        print(f"[3] VLM: SKIPPED (no images needed)", flush=True)


async def profile_vlm(images, api_key):
    import httpx
    from PIL import Image

    PROMPT = "Describe this image for a pharmaceutical knowledge base. Extract data points, describe charts, tables, or chemical structures."

    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    sem = asyncio.Semaphore(2)

    async def call(client, idx, img_bytes, label):
        async with sem:
            # Encode image
            img = Image.open(io.BytesIO(img_bytes))
            if img.mode != 'RGB':
                img = img.convert('RGB')
            max_dim = 1280
            if max(img.size) > max_dim:
                img.thumbnail((max_dim, max_dim), Image.Resampling.LANCZOS)
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=80)
            b64 = base64.b64encode(buf.getvalue()).decode('utf-8')

            payload = {
                "model": "mistral-small-latest",
                "messages": [{"role": "user", "content": [
                    {"type": "text", "text": PROMPT},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64}"}}
                ]}],
                "max_tokens": 1024, "temperature": 0.15
            }
            t = time.time()
            try:
                resp = await client.post("https://api.mistral.ai/v1/chat/completions",
                                        headers=headers, json=payload, timeout=60.0)
                e = time.time() - t
                if resp.status_code == 200:
                    c = resp.json()['choices'][0]['message']['content']
                    return label, e, len(c), "ok"
                elif resp.status_code == 429:
                    return label, e, 0, f"429 rate-limited"
                return label, e, 0, f"ERROR {resp.status_code}"
            except Exception as exc:
                return label, time.time()-t, 0, f"{type(exc).__name__}"

    t0 = time.time()
    async with httpx.AsyncClient() as client:
        tasks = [call(client, i[0], i[1], i[2]) for i in images]
        results = await asyncio.gather(*tasks)

    print(f"[3] VLM (Mistral small, parallel): {time.time()-t0:.2f}s", flush=True)
    for label, elapsed, chars, status in results:
        print(f"    {label}: {elapsed:.2f}s, {chars} chars, {status}", flush=True)


if __name__ == "__main__":
    asyncio.run(main())
