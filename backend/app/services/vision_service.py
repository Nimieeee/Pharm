"""
Hybrid Document Vision Service
- PDFs: Extract text with PyMuPDF, only send embedded images to VLM
- PPTX: Extract text with python-pptx, only send embedded images to VLM
- Images: Send directly to VLM
- Fallback: Full-page VLM for scanned/unreadable pages
"""
import base64
import io
import os
import subprocess
import tempfile
import logging
import asyncio
import time as _time
from typing import List, Any, Dict, Tuple, Optional

try:
    from PIL import Image
    from pdf2image import convert_from_bytes, convert_from_path
except ImportError:
    Image = None
    convert_from_bytes = None
    convert_from_path = None

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

import httpx

logger = logging.getLogger(__name__)

# --- CONFIGURATION ---
# Use Mistral's own API for vision — faster, simpler, already has API key
# Import settings to get API keys from .env properly
try:
    from app.core.config import settings
    MISTRAL_API_KEY = settings.MISTRAL_API_KEY or ""
    NVIDIA_API_KEY = settings.NVIDIA_API_KEY or ""
except Exception:
    # Fallback for direct module testing
    MISTRAL_API_KEY = os.getenv("MISTRAL_API_KEY", "")
    NVIDIA_API_KEY = os.getenv("NVIDIA_API_KEY", "")

MISTRAL_API_URL = "https://api.mistral.ai/v1/chat/completions"

# mistral-small-latest has vision support and is fast for OCR/image description
VISION_MODEL = "mistral-small-latest"
VISION_CONCURRENCY = 2
MIN_TEXT_CHARS = 100  # Pages with less text than this are treated as "scanned"
MIN_IMAGE_BYTES = 50000  # 50KB — skip tiny icons, logos, backgrounds
MIN_IMAGE_DIM = 200  # Skip images smaller than 200x200
MAX_IMAGES_PER_DOC = 20  # Cap total VLM calls per document

IMAGE_DESCRIPTION_PROMPT = """Describe this image in detail for a pharmaceutical knowledge base.
If it's a chart/graph: describe axes, data points, trends, and legends.
If it's a chemical structure: describe the molecular structure.
If it's a table: convert to Markdown table format.
If it's a diagram: describe all components and relationships."""

# ============================================================
# PUBLIC API — Called by smart_loader.py
# ============================================================

async def process_pdf_hybrid(
    content: bytes, filename: str, user_prompt: str, api_key: str,
    mode: str = "detailed", chunk_callback: Any = None
) -> str:
    """Hybrid PDF processing: text extraction + VLM only for images."""
    t0 = _time.time()
    logger.info(f"📄 [Hybrid PDF] Processing {filename}...")

    if not fitz:
        logger.warning("⚠️ PyMuPDF not available, falling back to full-vision pipeline")
        return await process_visual_document(content, filename, user_prompt, api_key, mode, chunk_callback)

    doc = fitz.open(stream=content, filetype="pdf")
    pages_content = []
    images_to_analyze = []  # List of (page_idx, image_bytes, context_label, xref)
    seen_xrefs = set()  # Deduplicate repeated images (templates)
    xref_page_count = {}  # Track how many pages each image appears on

    t_extract = _time.time()
    for page_idx in range(len(doc)):
        page = doc[page_idx]
        text = page.get_text("text").strip()

        # Extract embedded images from this page (deduplicate by xref)
        page_images = page.get_images(full=True)
        for img_idx, img_info in enumerate(page_images):
            try:
                xref = img_info[0]
                if xref in seen_xrefs:
                    # Track how many pages this image appears on (for template detection)
                    xref_page_count[xref] = xref_page_count.get(xref, 1) + 1
                    continue
                seen_xrefs.add(xref)
                xref_page_count[xref] = 1
                
                img_data = doc.extract_image(xref)
                if img_data and img_data.get("image"):
                    img_bytes = img_data["image"]
                    if len(img_bytes) > MIN_IMAGE_BYTES:
                        w = img_data.get("width", 0)
                        h = img_data.get("height", 0)
                        if w >= MIN_IMAGE_DIM and h >= MIN_IMAGE_DIM:
                            images_to_analyze.append((
                                page_idx,
                                img_bytes,
                                f"Page {page_idx+1}, Image {img_idx+1}",
                                xref
                            ))
            except Exception as e:
                logger.debug(f"Could not extract image from page {page_idx+1}: {e}")

        if len(text) >= MIN_TEXT_CHARS:
            # Text-rich page — use extracted text directly
            pages_content.append((page_idx, f"## Page {page_idx+1}\n{text}"))
            logger.debug(f"📝 Page {page_idx+1}: text extraction ({len(text)} chars)")
        else:
            # Scanned/image-only page — render to image for VLM
            logger.debug(f"👁️ Page {page_idx+1}: scanned page, queuing for VLM")
            try:
                pix = page.get_pixmap(dpi=150)
                img_bytes = pix.tobytes("jpeg")
                # Use None for xref since this is a full-page render, not an embedded image
                images_to_analyze.append((page_idx, img_bytes, f"Page {page_idx+1} (full page scan)", None))
            except Exception as e:
                logger.error(f"Could not render page {page_idx+1}: {e}")
                pages_content.append((page_idx, f"## Page {page_idx+1}\n[Could not extract content]"))

    # Filter out template images (appear on >50% of pages = likely decoration)
    total_pages = len(doc)
    doc.close()
    
    if images_to_analyze:
        before_filter = len(images_to_analyze)
        images_to_analyze = [
            img for img in images_to_analyze
            # Keep if: no xref (full page scan) OR xref appears on <50% of pages
            if img[3] is None or xref_page_count.get(img[3], 1) < total_pages * 0.5
        ]
        if len(images_to_analyze) < before_filter:
            logger.info(f"📊 Filtered {before_filter - len(images_to_analyze)} template images (appeared on >50% of pages)")
    
    t_extract_done = _time.time()
    # Cap images to avoid excessive VLM calls
    if len(images_to_analyze) > MAX_IMAGES_PER_DOC:
        logger.info(f"📊 Capping images from {len(images_to_analyze)} to {MAX_IMAGES_PER_DOC}")
        # Keep the largest images (most likely to be charts/diagrams)
        images_to_analyze.sort(key=lambda x: len(x[1]), reverse=True)
        images_to_analyze = images_to_analyze[:MAX_IMAGES_PER_DOC]

    logger.info(f"⏱️ [Hybrid PDF] Text extraction: {t_extract_done - t_extract:.1f}s, "
                f"{len(pages_content)} text pages, {len(images_to_analyze)} images queued")

    # Send only queued images to VLM
    image_descriptions = {}
    if images_to_analyze:
        image_descriptions = await _batch_vision_analyze(images_to_analyze)
        logger.info(f"📊 [Hybrid PDF] VLM returned {len(image_descriptions)} image descriptions")

    # Combine text pages + image descriptions, sorted by page order
    all_content = list(pages_content)
    for page_idx, desc in image_descriptions.items():
        all_content.append((page_idx, desc))

    all_content.sort(key=lambda x: x[0])
    full_text = "\n\n".join([c[1] for c in all_content])

    # Log warning if no content extracted
    if not full_text.strip():
        logger.error(f"❌ [Hybrid PDF] {filename}: No content extracted! "
                    f"Text pages: {len(pages_content)}, Images analyzed: {len(images_to_analyze)}, "
                    f"Image descriptions: {len(image_descriptions)}")
        # Return a placeholder to avoid 0 chunks
        return f"[PDF Processing Error: Could not extract content from {filename}. " \
               f"This may be a scanned/image-based PDF that requires OCR. " \
               f"Pages: {total_pages}, Text pages: {len(pages_content)}, Images: {len(images_to_analyze)}]"

    t_done = _time.time()
    logger.info(f"⏱️ [Hybrid PDF] Total: {t_done - t0:.1f}s (extract={t_extract_done - t_extract:.1f}s, "
                f"vision={t_done - t_extract_done:.1f}s for {len(images_to_analyze)} images), "
                f"extracted {len(full_text)} chars")

    return full_text


async def process_pptx_hybrid(
    content: bytes, filename: str, user_prompt: str, api_key: str,
    mode: str = "detailed", chunk_callback: Any = None
) -> str:
    """Hybrid PPTX processing: text from python-pptx + VLM for embedded images."""
    t0 = _time.time()
    logger.info(f"📄 [Hybrid PPTX] Processing {filename}...")

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        logger.warning("⚠️ python-pptx not available, falling back to full-vision pipeline")
        return await process_visual_document(content, filename, user_prompt, api_key, mode, chunk_callback)

    prs = Presentation(io.BytesIO(content))
    slides_content = []
    images_to_analyze = []

    t_extract = _time.time()
    for slide_idx, slide in enumerate(prs.slides):
        texts = []

        for shape in slide.shapes:
            # Extract text from text frames
            if shape.has_text_frame:
                frame_text = shape.text_frame.text.strip()
                if frame_text:
                    texts.append(frame_text)

            # Extract tables as markdown
            if shape.has_table:
                texts.append(_table_to_markdown(shape.table))

            # Queue embedded images for VLM
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_bytes = shape.image.blob
                    if len(img_bytes) > MIN_IMAGE_BYTES:  # Skip tiny icons
                        images_to_analyze.append((
                            slide_idx,
                            img_bytes,
                            f"Slide {slide_idx+1}, Image"
                        ))
                except Exception as e:
                    logger.debug(f"Could not extract image from slide {slide_idx+1}: {e}")

            # Handle grouped shapes
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    for sub_shape in shape.shapes:
                        if hasattr(sub_shape, 'text_frame') and sub_shape.has_text_frame:
                            t = sub_shape.text_frame.text.strip()
                            if t:
                                texts.append(t)
                except Exception:
                    pass

        slide_text = "\n".join(texts)
        if slide_text.strip():
            slides_content.append((slide_idx, f"## Slide {slide_idx+1}\n{slide_text}"))
        else:
            # Empty text slide — render it for VLM (might be a diagram-only slide)
            slides_content.append((slide_idx, f"## Slide {slide_idx+1}\n[Visual content - see image descriptions]"))

    t_extract_done = _time.time()
    logger.info(f"⏱️ [Hybrid PPTX] Text extraction: {t_extract_done - t_extract:.1f}s, "
                f"{len(slides_content)} slides, {len(images_to_analyze)} images queued")

    # Send images to VLM
    image_descriptions = {}
    if images_to_analyze:
        image_descriptions = await _batch_vision_analyze(images_to_analyze)

    # Combine slide text + image descriptions
    all_content = list(slides_content)
    for slide_idx, desc in image_descriptions.items():
        # Append image descriptions to the corresponding slide
        for i, (idx, text) in enumerate(all_content):
            if idx == slide_idx:
                all_content[i] = (idx, text + "\n\n" + desc)
                break
        else:
            all_content.append((slide_idx, desc))

    all_content.sort(key=lambda x: x[0])
    full_text = "\n\n".join([c[1] for c in all_content])

    t_done = _time.time()
    logger.info(f"⏱️ [Hybrid PPTX] Total: {t_done - t0:.1f}s (extract={t_extract_done - t_extract:.1f}s, "
                f"vision={t_done - t_extract_done:.1f}s for {len(images_to_analyze)} images)")

    return full_text


async def process_visual_document(
    content: bytes, filename: str, user_prompt: str, api_key: str,
    mode: str = "detailed", chunk_callback: Any = None
) -> str:
    """Full-vision pipeline for standalone images. Also used as fallback."""
    t0 = _time.time()
    logger.info(f"🖼️ [Vision] Processing image {filename}...")

    if not Image:
        return f"Error: Pillow not installed"

    # For standalone images, just send directly
    ext = filename.split('.')[-1].lower()
    if ext in ['png', 'jpg', 'jpeg', 'webp', 'gif', 'bmp']:
        img = Image.open(io.BytesIO(content))
        b64 = _optimize_and_encode(img)
        results = await _batch_vision_analyze([(0, None, "Image")], pre_encoded=[b64])
        if results:
            return list(results.values())[0]
        return "Error: Could not analyze image"

    # For PDF/PPTX fallback (when PyMuPDF/python-pptx unavailable)
    images = _convert_to_images(content, filename)
    if not images:
        return f"Error: Could not extract images from {filename}"

    encoded = [_optimize_and_encode(img) for img in images]
    items = [(i, None, f"Page {i+1}") for i in range(len(encoded))]
    results = await _batch_vision_analyze(items, pre_encoded=encoded)

    sorted_results = sorted(results.items(), key=lambda x: x[0])
    full_text = "\n\n".join([v for _, v in sorted_results])

    t_done = _time.time()
    logger.info(f"⏱️ [Vision] Total: {t_done - t0:.1f}s for {len(images)} pages")
    return full_text


# ============================================================
# SHARED VLM CALLER
# ============================================================

async def _batch_vision_analyze(
    items: List[Tuple[int, Optional[bytes], str]],
    pre_encoded: List[str] = None,
    api_key: str = None
) -> Dict[int, str]:
    """
    Analyze multiple images via VLM with concurrency control.

    Args:
        items: List of (page_idx, image_bytes, label). image_bytes can be None if pre_encoded.
        pre_encoded: Optional list of pre-encoded base64 strings (same length as items).
        api_key: Optional API key (uses global MISTRAL_API_KEY if not provided).

    Returns:
        Dict mapping page_idx -> description text
    """
    if not items:
        return {}

    # Use provided API key or fall back to global
    effective_api_key = api_key or MISTRAL_API_KEY

    semaphore = asyncio.Semaphore(VISION_CONCURRENCY)
    results = {}

    if not effective_api_key:
        logger.error("❌ MISTRAL_API_KEY not set — cannot analyze images")
        return {}

    headers = {
        "Authorization": f"Bearer {effective_api_key}",
        "Content-Type": "application/json",
        "Accept": "application/json"
    }

    async def analyze_single(idx, item_idx, label, http_client):
        b64 = None
        if pre_encoded and item_idx < len(pre_encoded):
            b64 = pre_encoded[item_idx]
        else:
            img_bytes = items[item_idx][1]
            if img_bytes and Image:
                img = Image.open(io.BytesIO(img_bytes))
                b64 = _optimize_and_encode(img)

        if not b64:
            return idx, f"[Could not encode image for {label}]"

        async with semaphore:
            max_retries = 3
            for attempt in range(max_retries):
                try:
                    payload = {
                        "model": VISION_MODEL,
                        "messages": [{
                            "role": "user",
                            "content": [
                                {"type": "text", "text": IMAGE_DESCRIPTION_PROMPT},
                                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{b64}"}}
                            ]
                        }],
                        "max_tokens": 1024,
                        "temperature": 0.15
                    }

                    resp = await http_client.post(MISTRAL_API_URL, headers=headers, json=payload, timeout=60.0)

                    if resp.status_code == 200:
                        data = resp.json()
                        desc = data['choices'][0]['message']['content']
                        logger.info(f"✅ {label}: {len(desc)} chars")
                        return idx, f"### {label}\n{desc}"

                    elif resp.status_code == 429:
                        wait = 2 ** (attempt + 1)
                        logger.warning(f"⚠️ {label} rate-limited, retry in {wait}s")
                        await asyncio.sleep(wait)
                        continue
                    else:
                        logger.error(f"❌ {label}: API error {resp.status_code}")
                        return idx, f"[{label}: API Error {resp.status_code}]"

                except Exception as e:
                    if attempt < max_retries - 1:
                        await asyncio.sleep(2 ** (attempt + 1))
                        continue
                    logger.error(f"❌ {label}: {e}")
                    return idx, f"[{label}: Error - {str(e)[:100]}]"

            return idx, f"[{label}: Max retries exceeded]"

    t0 = _time.time()
    logger.info(f"👁️ [VLM] Sending {len(items)} images to {VISION_MODEL} (concurrency={VISION_CONCURRENCY})...")

    async with httpx.AsyncClient() as http_client:
        tasks = [
            analyze_single(items[i][0], i, items[i][2], http_client)
            for i in range(len(items))
        ]
        raw_results = await asyncio.gather(*tasks)

    for page_idx, desc in raw_results:
        if page_idx in results:
            results[page_idx] += "\n\n" + desc
        else:
            results[page_idx] = desc

    logger.info(f"⏱️ [VLM] {len(items)} images analyzed in {_time.time()-t0:.1f}s")
    return results


# ============================================================
# HELPERS
# ============================================================

def _optimize_and_encode(image) -> str:
    """Resize and encode image to base64 JPEG."""
    if not Image:
        return ""

    max_dim = 1280
    img_copy = image.copy()

    if img_copy.mode == 'RGBA':
        bg = Image.new('RGB', img_copy.size, (255, 255, 255))
        bg.paste(img_copy, mask=img_copy.split()[3])
        img_copy = bg
    elif img_copy.mode != 'RGB':
        img_copy = img_copy.convert('RGB')

    if max(img_copy.size) > max_dim:
        img_copy.thumbnail((max_dim, max_dim), Image.Resampling.LANCZOS)

    buf = io.BytesIO()
    img_copy.save(buf, format="JPEG", quality=80)
    return base64.b64encode(buf.getvalue()).decode('utf-8')


def _convert_to_images(content: bytes, filename: str):
    """Legacy converter for fallback pipeline."""
    ext = filename.split('.')[-1].lower()
    try:
        if ext == 'pdf':
            if convert_from_bytes:
                return convert_from_bytes(content, dpi=200, fmt="jpeg", thread_count=4)
            raise ImportError("pdf2image missing")
        elif ext == 'pptx':
            return _convert_pptx_robust(content)
        else:
            if Image:
                return [Image.open(io.BytesIO(content))]
            raise ImportError("Pillow missing")
    except Exception as e:
        logger.error(f"Image conversion failed: {e}")
        return []


def _convert_pptx_robust(content: bytes):
    """PPTX → PDF → images (legacy fallback)."""
    if not convert_from_path:
        raise ImportError("pdf2image required for pptx")

    with tempfile.TemporaryDirectory() as td:
        inp = os.path.join(td, "temp.pptx")
        with open(inp, "wb") as f:
            f.write(content)
        try:
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf", inp, "--outdir", td],
                stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, check=True
            )
        except Exception:
            logger.error("LibreOffice conversion failed")
            return []
        pdf_path = os.path.join(td, "temp.pdf")
        if not os.path.exists(pdf_path):
            return []
        return convert_from_path(pdf_path, dpi=200)


def _table_to_markdown(table) -> str:
    """Convert a python-pptx Table object to Markdown."""
    rows = []
    for row_idx, row in enumerate(table.rows):
        cells = [cell.text.strip() for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
        if row_idx == 0:
            rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
    return "\n".join(rows)


# ============================================================
# VISION SERVICE CLASS (for ServiceContainer registration)
# ============================================================

class VisionService:
    """
    Vision service for image analysis using Mistral VLM.

    Following CLAUDE.md patterns:
    - Uses ServiceContainer for dependencies
    - Lazy loading for API keys
    - Centralized image processing logic
    """

    def __init__(self):
        self._container = None
        self._api_key = None

    @property
    def container(self):
        """Get container - should be initialized at app startup"""
        if self._container is None:
            from app.core.container import container
            self._container = container
        return self._container

    @property
    def api_key(self) -> str:
        """Get API key from settings"""
        if self._api_key is None:
            from app.core.config import settings
            self._api_key = settings.MISTRAL_API_KEY or ""
        return self._api_key

    async def analyze_image(self, image_url: str, prompt: str = None) -> str:
        """
        Analyze a single image from URL.

        Args:
            image_url: Public URL of the image
            prompt: Optional custom prompt (uses default if not provided)

        Returns:
            Image description text
        """
        if not self.api_key:
            logger.error("❌ MISTRAL_API_KEY not set — cannot analyze images")
            return "[Image analysis unavailable - API key not configured]"

        try:
            # Download image and convert to base64
            async with httpx.AsyncClient() as client:
                resp = await client.get(image_url, timeout=30.0)
                if resp.status_code != 200:
                    return f"[Failed to download image: HTTP {resp.status_code}]"

                image_bytes = resp.content

            # Encode image
            if Image:
                img = Image.open(io.BytesIO(image_bytes))
                b64 = _optimize_and_encode(img)
            else:
                return "[Pillow not installed - cannot process images]"

            if not b64:
                return "[Could not encode image]"

            # Call VLM
            results = await _batch_vision_analyze(
                items=[(0, None, "Image")],
                pre_encoded=[b64],
                api_key=self.api_key
            )

            return results.get(0, "[No analysis result]")

        except Exception as e:
            logger.error(f"❌ Image analysis error: {e}")
            return f"[Image analysis error: {str(e)[:200]}]"

    async def analyze_image_bytes(self, image_bytes: bytes, label: str = "Image") -> str:
        """
        Analyze image from raw bytes.

        Args:
            image_bytes: Raw image bytes
            label: Label for logging

        Returns:
            Image description text
        """
        if not self.api_key:
            logger.error("❌ MISTRAL_API_KEY not set — cannot analyze images")
            return "[Image analysis unavailable - API key not configured]"

        try:
            # Encode image
            if Image:
                img = Image.open(io.BytesIO(image_bytes))
                b64 = _optimize_and_encode(img)
            else:
                return "[Pillow not installed - cannot process images]"

            if not b64:
                return "[Could not encode image]"

            # Call VLM
            results = await _batch_vision_analyze(
                items=[(0, None, label)],
                pre_encoded=[b64],
                api_key=self.api_key
            )

            return results.get(0, "[No analysis result]")

        except Exception as e:
            logger.error(f"❌ Image analysis error: {e}")
            return f"[Image analysis error: {str(e)[:200]}]"


# Singleton instance for backward compatibility
_vision_service: Optional[VisionService] = None

def get_vision_service() -> VisionService:
    """Get or create singleton VisionService"""
    global _vision_service
    if _vision_service is None:
        _vision_service = VisionService()
    return _vision_service
