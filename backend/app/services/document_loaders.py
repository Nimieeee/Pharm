"""
Enhanced Document Loaders using LangChain
Handles document loading and processing with comprehensive error handling
"""

import os
import tempfile
import logging
import time
from typing import List, Dict, Any, Optional
from pathlib import Path

from langchain_core.documents import Document
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
    Docx2txtLoader,
    UnstructuredPowerPointLoader,
)
import pandas as pd

from app.core.config import settings
from app.core.logging_config import RAGLogger

logger = logging.getLogger(__name__)
rag_logger = RAGLogger(__name__)


class DocumentProcessingError(Exception):
    """
    Custom exception for document processing errors
    
    Attributes:
        message: Human-readable error message
        error_category: Categorized error type for programmatic handling
        is_user_error: True if error is due to user input, False for system errors
        details: Additional error context
    """
    def __init__(
        self, 
        message: str, 
        error_category: str = "processing_error",
        is_user_error: bool = True,
        details: Optional[Dict[str, Any]] = None
    ):
        self.message = message
        self.error_category = error_category
        self.is_user_error = is_user_error
        self.details = details or {}
        super().__init__(self.message)
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert exception to dictionary for API responses"""
        return {
            "error": self.message,
            "error_category": self.error_category,
            "error_type": "user_error" if self.is_user_error else "system_error",
            "details": self.details
        }


class ErrorCategory:
    """Error category constants for consistent error handling"""
    EMPTY_CONTENT = "empty_content"
    ENCODING_ERROR = "encoding_error"
    CORRUPTED_FILE = "corrupted_file"
    UNSUPPORTED_FORMAT = "unsupported_format"
    INSUFFICIENT_CONTENT = "insufficient_content"
    PROCESSING_ERROR = "processing_error"
    VALIDATION_ERROR = "validation_error"
    FILE_NOT_FOUND = "file_not_found"
    PERMISSION_ERROR = "permission_error"


class ErrorMessageTemplates:
    """Standardized error message templates for consistent user feedback"""
    
    @staticmethod
    def unsupported_format(filename: str, extension: str, supported_formats: List[str]) -> str:
        """Error message for unsupported file formats"""
        formats_list = ", ".join(supported_formats)
        return (
            f"Unsupported file format: {extension}. "
            f"File '{filename}' cannot be processed. "
            f"Supported formats: {formats_list}"
        )
    
    @staticmethod
    def empty_content(filename: str) -> str:
        """Error message for files with no extractable content"""
        return (
            f"No content could be extracted from '{filename}'. "
            f"The file appears to be empty or contains only non-text elements."
        )
    
    @staticmethod
    def encoding_error(filename: str, attempted_encodings: List[str]) -> str:
        """Error message for text encoding failures"""
        encodings_list = ", ".join(attempted_encodings)
        return (
            f"Could not decode text file '{filename}'. "
            f"Attempted encodings: {encodings_list}. "
            f"Please ensure the file is saved in UTF-8, UTF-16, or Latin-1 encoding."
        )
    
    @staticmethod
    def corrupted_file(filename: str, file_type: str) -> str:
        """Error message for corrupted or invalid files"""
        return (
            f"The file '{filename}' appears to be corrupted or invalid. "
            f"Unable to parse as {file_type} format. "
            f"Please try re-saving or re-exporting the file."
        )
    
    @staticmethod
    def insufficient_content(filename: str, char_count: int, min_required: int) -> str:
        """Error message for files with insufficient content"""
        return (
            f"Insufficient content in '{filename}'. "
            f"Found {char_count} characters, minimum required is {min_required}. "
            f"Please provide a file with more substantial content."
        )
    
    @staticmethod
    def processing_error(filename: str, error_details: str) -> str:
        """Generic processing error message"""
        return (
            f"Failed to process '{filename}'. "
            f"Error: {error_details}"
        )
    
    @staticmethod
    def validation_error(filename: str, validation_issue: str) -> str:
        """Error message for validation failures"""
        return (
            f"Validation failed for '{filename}'. "
            f"Issue: {validation_issue}"
        )
    
    @staticmethod
    def file_not_found(filename: str) -> str:
        """Error message for missing files"""
        return (
            f"File '{filename}' not found or cannot be accessed."
        )
    
    @staticmethod
    def permission_error(filename: str) -> str:
        """Error message for permission issues"""
        return (
            f"Permission denied when accessing '{filename}'. "
            f"Please check file permissions."
        )


class EnhancedDocumentLoader:
    """Enhanced document loader using LangChain with comprehensive error handling"""
    
    def __init__(self):
        self.supported_extensions = {
            '.pdf': self._load_pdf,
            '.txt': self._load_text,
            '.md': self._load_text,
            '.docx': self._load_docx,
            '.pptx': self._load_pptx,
            '.xlsx': self._load_xlsx,
            '.csv': self._load_csv,
            '.png': self._load_image,
            '.jpg': self._load_image,
            '.jpeg': self._load_image,
            '.gif': self._load_image,
            '.bmp': self._load_image,
            '.webp': self._load_image,
            '.sdf': self._load_sdf,
            '.mol': self._load_sdf,
        }
        logger.info("✅ Enhanced document loader initialized (supports text, PDF, DOCX, PPTX, XLSX, CSV, SDF, and images)")
    
    def is_supported_format(self, filename: str) -> bool:
        """Check if file format is supported"""
        file_extension = Path(filename).suffix.lower()
        return file_extension in self.supported_extensions
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats"""
        return list(self.supported_extensions.keys())
    
    async def load_document(
        self, 
        file_content: bytes, 
        filename: str, 
        additional_metadata: Optional[Dict[str, Any]] = None,
        image_analyzer: Optional[Any] = None,
        user_prompt: Optional[str] = None,
        mode: str = "detailed",
        chunk_callback: Optional[Any] = None
    ) -> List[Document]:
        """
        Smart Vision Router Implementation:
        Delegates to smart_loader for intelligent document processing.
        Falls back to legacy LangChain loaders if smart processing fails.
        """
        try:
            start_time = time.time()
            from app.services.smart_loader import process_file as smart_process
            from app.core.config import settings

            # Basic Validation
            if not file_content:
                 raise DocumentProcessingError(
                    ErrorMessageTemplates.empty_content(filename), 
                    error_category=ErrorCategory.EMPTY_CONTENT,
                    is_user_error=True
                 )

            if not self.is_supported_format(filename):
                ext = Path(filename).suffix.lower()
                raise DocumentProcessingError(
                    ErrorMessageTemplates.unsupported_format(filename, ext, self.get_supported_formats()),
                    error_category=ErrorCategory.UNSUPPORTED_FORMAT
                )

            # Default Prompt Strategy
            if not user_prompt:
                user_prompt = (
                    "Analyze this pharmaceutical document comprehensively. "
                    "Extract all text verbatim, interpret charts/tables data, "
                    "and provide statistical summary for data files."
                )

            logger.info(f"🧠 Smart Router: Processing {filename} in mode={mode}...")

            # Process using Smart Vision Router
            content = await smart_process(
                file_content=file_content,
                filename=filename,
                user_prompt=user_prompt,
                api_key=settings.MISTRAL_API_KEY,
                mode=mode,
                chunk_callback=chunk_callback
            )

            # Check for explicit failure strings from smart_loader
            if content.startswith("Error") or content.startswith("❌"):
                 # Force fallback
                 raise Exception(f"Smart Loader returned error: {content}")

            # Success - Create Document
            metadata = additional_metadata or {}
            ext = Path(filename).suffix.lower()
            metadata.update({
                "source": filename,
                "file_extension": ext,
                "processor": "smart_vision_router",
                "processed_at": time.time(),
                "filename": filename,
                "loader_version": "smart_v2",
                "file_type": ext.lstrip('.')
            })
            
            logger.info(f"✅ Smart Vision Router processed {filename} in {time.time()-start_time:.2f}s")
            return [Document(page_content=content, metadata=metadata)]

        except DocumentProcessingError as dpe:
            # If it's a validation error (empty content, unsupported format), re-raise
            # Don't fallback for user errors
            logger.warning(f"Smart Loader Validation Error: {dpe.message}")
            if dpe.is_user_error:
                raise dpe
            
            # If system error, try fallback
            logger.info("⚠️ Falling back to legacy loader due to processing error...")
            return await self._legacy_load_document(file_content, filename, additional_metadata, image_analyzer)

        except Exception as e:
            logger.error(f"Smart Loader System Failure: {e}")
            logger.info("⚠️ Falling back to legacy loader...")
            return await self._legacy_load_document(file_content, filename, additional_metadata, image_analyzer)

    async def _legacy_load_document(
        self, 
        file_content: bytes, 
        filename: str, 
        additional_metadata: Optional[Dict[str, Any]] = None,
        image_analyzer: Optional[Any] = None  # Callable[[bytes], Awaitable[str]]
    ) -> List[Document]:
        """
        Load document from file content using appropriate LangChain loader
        
        Args:
            file_content: Raw file bytes
            filename: Original filename
            additional_metadata: Additional metadata to add to documents
            image_analyzer: Optional async function to analyze images found in documents
            
        Returns:
            List of LangChain Document objects
            
        Raises:
            DocumentProcessingError: If document processing fails
        """
        start_time = time.time()
        file_size = len(file_content) if file_content else 0
        file_extension = Path(filename).suffix.lower()
        
        # Log document loading start
        logger.info(
            f"📄 Starting document load: {filename} ({file_size} bytes, {file_extension})",
            extra={
                'operation': 'document_load',
                'document_name': filename,
                'file_size': file_size,
                'file_type': file_extension
            }
        )
        
        if not file_content:
            error_msg = ErrorMessageTemplates.empty_content(filename)
            logger.error(
                f"❌ Document load failed: {filename} - Empty file",
                extra={
                    'operation': 'document_load',
                    'document_name': filename,
                    'file_size': 0,
                    'error_type': ErrorCategory.EMPTY_CONTENT,
                    'duration': (time.time() - start_time) * 1000
                }
            )
            raise DocumentProcessingError(
                error_msg,
                error_category=ErrorCategory.EMPTY_CONTENT,
                is_user_error=True,
                details={"filename": filename, "file_size": 0}
            )
        
        if not self.is_supported_format(filename):
            error_msg = ErrorMessageTemplates.unsupported_format(
                filename, 
                file_extension, 
                self.get_supported_formats()
            )
            logger.error(
                f"❌ Document load failed: {filename} - Unsupported format {file_extension}",
                extra={
                    'operation': 'document_load',
                    'document_name': filename,
                    'file_type': file_extension,
                    'error_type': ErrorCategory.UNSUPPORTED_FORMAT,
                    'duration': (time.time() - start_time) * 1000
                }
            )
            raise DocumentProcessingError(
                error_msg,
                error_category=ErrorCategory.UNSUPPORTED_FORMAT,
                is_user_error=True,
                details={
                    "filename": filename,
                    "extension": file_extension,
                    "supported_formats": self.get_supported_formats()
                }
            )
        
        # Create temporary file
        tmp_file_path = None
        tmp_file_creation_time = time.time()
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_file:
                tmp_file.write(file_content)
                tmp_file_path = tmp_file.name
            
            tmp_file_duration = time.time() - tmp_file_creation_time
            logger.debug(
                f"📝 Created temporary file for {filename} ({tmp_file_duration:.3f}s)",
                extra={
                    'operation': 'temp_file_creation',
                    'document_name': filename,
                    'duration': tmp_file_duration * 1000
                }
            )
        except Exception as e:
            logger.error(
                f"❌ Failed to create temporary file for {filename}: {e}",
                extra={
                    'operation': 'temp_file_creation',
                    'document_name': filename,
                    'error_type': 'temp_file_error',
                    'duration': (time.time() - start_time) * 1000
                },
                exc_info=True
            )
            raise DocumentProcessingError(
                ErrorMessageTemplates.processing_error(filename, "Failed to create temporary file"),
                error_category=ErrorCategory.PROCESSING_ERROR,
                is_user_error=False,
                details={"filename": filename, "error": str(e)}
            )
        
        try:
            # Load document using appropriate loader
            loader_start_time = time.time()
            loader_func = self.supported_extensions[file_extension]
            loader_name = loader_func.__name__
            
            logger.debug(
                f"🔧 Using loader {loader_name} for {filename}",
                extra={
                    'operation': 'loader_selection',
                    'document_name': filename,
                    'loader': loader_name,
                    'file_type': file_extension
                }
            )
            
            # Pass image_analyzer to appropriate loaders
            if file_extension == '.pdf':
                documents = await self._load_pdf(tmp_file_path, filename, image_analyzer)
            elif file_extension == '.docx':
                documents = await self._load_docx(tmp_file_path, filename, image_analyzer)
            elif file_extension == '.pptx':
                documents = await self._load_pptx(tmp_file_path, filename, image_analyzer)
            else:
                documents = await loader_func(tmp_file_path, filename)
                
            loader_duration = time.time() - loader_start_time
            
            if not documents:
                logger.error(
                    f"❌ Document load failed: {filename} - No documents extracted",
                    extra={
                        'operation': 'document_load',
                        'document_name': filename,
                        'file_type': file_extension,
                        'loader': loader_name,
                        'error_type': ErrorCategory.EMPTY_CONTENT,
                        'duration': (time.time() - start_time) * 1000
                    }
                )
                raise DocumentProcessingError(
                    ErrorMessageTemplates.empty_content(filename),
                    error_category=ErrorCategory.EMPTY_CONTENT,
                    is_user_error=True,
                    details={"filename": filename, "file_type": file_extension}
                )
            
            # Add additional metadata
            if additional_metadata:
                for doc in documents:
                    doc.metadata.update(additional_metadata)
            
            # Add common metadata
            for doc in documents:
                doc.metadata.update({
                    "filename": filename,
                    "file_type": file_extension.lstrip('.'),
                    "loader_version": "langchain",
                    "total_pages": len(documents) if file_extension == '.pdf' else 1
                })
            
            # Calculate content statistics
            total_chars = sum(len(doc.page_content) for doc in documents)
            total_words = sum(len(doc.page_content.split()) for doc in documents)
            
            total_duration = time.time() - start_time
            
            # Log successful document load with comprehensive statistics
            logger.info(
                f"✅ Successfully loaded {len(documents)} document(s) from {filename} "
                f"({total_chars} chars, {total_words} words, {total_duration:.2f}s)",
                extra={
                    'operation': 'document_load',
                    'document_name': filename,
                    'file_size': file_size,
                    'file_type': file_extension,
                    'loader': loader_name,
                    'document_count': len(documents),
                    'total_characters': total_chars,
                    'total_words': total_words,
                    'loader_duration': loader_duration * 1000,
                    'duration': total_duration * 1000
                }
            )
            
            return documents
            
        except DocumentProcessingError:
            raise
        except Exception as e:
            total_duration = time.time() - start_time
            logger.error(
                f"❌ Error loading document {filename}: {e}",
                extra={
                    'operation': 'document_load',
                    'document_name': filename,
                    'file_type': file_extension,
                    'error_type': ErrorCategory.PROCESSING_ERROR,
                    'duration': total_duration * 1000
                },
                exc_info=True
            )
            raise DocumentProcessingError(
                ErrorMessageTemplates.processing_error(filename, str(e)),
                error_category=ErrorCategory.PROCESSING_ERROR,
                is_user_error=False,
                details={"filename": filename, "error": str(e), "file_type": file_extension}
            )
        finally:
            # Clean up temporary file
            if tmp_file_path:
                try:
                    os.unlink(tmp_file_path)
                    logger.debug(
                        f"🗑️  Cleaned up temporary file for {filename}",
                        extra={'operation': 'temp_file_cleanup', 'document_name': filename}
                    )
                except Exception as e:
                    logger.warning(
                        f"⚠️  Failed to clean up temporary file {tmp_file_path}: {e}",
                        extra={
                            'operation': 'temp_file_cleanup',
                            'document_name': filename,
                            'error_type': 'cleanup_error'
                        }
                    )
    
    async def _load_pdf(
        self, 
        file_path: str, 
        filename: str,
        image_analyzer: Optional[Any] = None
    ) -> List[Document]:
        """
        Hybrid PDF Loader:
        1. Tries standard text extraction (PyPDF) first (Fast/Free).
        2. If text content is low (likely scanned), falls back to Vision Analysis (Pixtral).
        3. Allows manual 'image_analyzer' hook for Pixtral processing with Rate Limiting.
        """
        import time
        from app.utils.rate_limiter import mistral_limiter
        
        try:
            # 1. Try standard text extraction
            loader = PyPDFLoader(file_path)
            documents = loader.load()
            
            # Check text density
            total_text = sum(len(doc.page_content.strip()) for doc in documents)
            avg_text_per_page = total_text / len(documents) if documents else 0
            
            is_scanned = avg_text_per_page < 500
            
            if not is_scanned:
                logger.info(f"📄 PDF {filename} has sufficient text ({avg_text_per_page:.0f} chars/page). Using standard extraction.")
                # Add page numbers
                for i, doc in enumerate(documents):
                    doc.metadata.update({"page": i + 1, "source": filename, "loader": "PyPDFLoader"})
                    doc.page_content = self._clean_text(doc.page_content)
                return documents
            
            # 2. Fallback to Vision Analysis (Pixtral) if scanned
            logger.info(f"📷 PDF {filename} appears scanned/image-heavy ({avg_text_per_page:.0f} chars/page). Switching to Vision Analysis.")
            
            if not image_analyzer:
                logger.warning("⚠️ Scanned PDF detected but no image analyzer provided. Returning sparse text.")
                return documents

            # Import tools for converting PDF to images
            try:
                from pdf2image import convert_from_path
                import io
                import base64
            except ImportError as e:
                logger.error(f"Missing dependencies for PDF image conversion: {e}")
                raise DocumentProcessingError(
                    "Server missing dependencies (pdf2image/poppler) for scanned PDF processing.",
                    error_category=ErrorCategory.PROCESSING_ERROR
                )

            # Convert PDF to images
            try:
                images = convert_from_path(file_path)
            except Exception as e:
                logger.error(f"Failed to convert PDF to images: {e}")
                # Fallback to sparse text
                return documents
            
            vision_documents = []
            total_pages = len(images)
            
            import os
            api_key = os.getenv("MISTRAL_API_KEY")
            if not api_key:
                logger.warning("No Mistral API key for Vision analysis")
                return documents
                
            import httpx # Use httpx directly to avoid mistralai SDK dependency issues if not present
            
            for i, image in enumerate(images):
                # Rate Limit Wait
                await mistral_limiter.wait_for_slot()
                
                logger.info(f"👁️ Analyzing page {i+1}/{total_pages} of {filename} with Pixtral...")
                
                try:
                    # Prepare image
                    buffered = io.BytesIO()
                    image.save(buffered, format="JPEG")
                    base64_image = base64.b64encode(buffered.getvalue()).decode('utf-8')
                    
                    # Call Pixtral via HTTP directly (matches ai.py implementation)
                    async with httpx.AsyncClient(timeout=60.0) as client:
                        response = await client.post(
                            "https://api.mistral.ai/v1/chat/completions",
                            headers={
                                "Authorization": f"Bearer {api_key}",
                                "Content-Type": "application/json"
                            },
                            json={
                                "model": "pixtral-12b-2409",
                                "messages": [
                                    {
                                        "role": "user",
                                        "content": [
                                            {"type": "text", "text": "Analyze this page in detail. Transcribe all visible text. Describe charts, graphs, or chemical structures using Markdown."},
                                            {"type": "image_url", "image_url": f"data:image/jpeg;base64,{base64_image}"} 
                                        ]
                                    }
                                ],
                                "max_tokens": 2000
                            }
                        )
                        
                        if response.status_code == 200:
                            content = response.json()["choices"][0]["message"]["content"]
                            page_content = f"--- [Page {i+1} : Vision Analysis] ---\n{content}"
                        else:
                            logger.error(f"Vision API Error {response.status_code}: {response.text}")
                            page_content = f"--- [Page {i+1}] ---\n(Analysis failed: {response.status_code})"
                            
                    # Create document for this page
                    doc = Document(
                        page_content=page_content,
                        metadata={
                            "source": filename,
                            "page": i + 1,
                            "loader": "PixtralVisionLoader",
                            "is_scanned": True
                        }
                    )
                    vision_documents.append(doc)
                    
                except Exception as e:
                    logger.error(f"Error processing page {i+1}: {e}")
                    # Keep going for other pages
            
            if vision_documents:
                return vision_documents
            else:
                return documents # Fallback if vision failed completely
                
        except DocumentProcessingError:
            raise
        except Exception as e:
            logger.error(f"❌ PDF processing error for {filename}: {e}", exc_info=True)
            raise DocumentProcessingError(
                ErrorMessageTemplates.processing_error(filename, str(e)),
                error_category=ErrorCategory.PROCESSING_ERROR
            )
    
    async def _load_text(self, file_path: str, filename: str) -> List[Document]:
        """
        Load text document using TextLoader with enhanced encoding detection
        
        Tries multiple encodings in order of likelihood and logs which encoding succeeds
        """
        # Extended list of encodings to try, ordered by likelihood
        encodings_to_try = [
            'utf-8',           # Most common modern encoding
            'utf-8-sig',       # UTF-8 with BOM
            'utf-16',          # Common for Windows files
            'utf-16-le',       # UTF-16 Little Endian
            'utf-16-be',       # UTF-16 Big Endian
            'latin-1',         # ISO-8859-1, common for Western European
            'cp1252',          # Windows Western European
            'iso-8859-1',      # Latin-1 alias
            'ascii',           # Basic ASCII
            'cp437',           # DOS encoding
            'mac_roman',       # Classic Mac encoding
        ]
        
        last_error = None
        
        for encoding in encodings_to_try:
            try:
                loader = TextLoader(file_path, encoding=encoding)
                documents = loader.load()
                
                if documents and documents[0].page_content.strip():
                    # Clean and enhance the document
                    doc = documents[0]
                    doc.page_content = self._clean_text(doc.page_content)
                    doc.metadata.update({
                        "source": filename,
                        "encoding": encoding,
                        "loader": "TextLoader"
                    })
                    
                    logger.info(f"✅ Successfully loaded text file {filename} with encoding: {encoding}")
                    return [doc]
                    
            except UnicodeDecodeError as e:
                logger.debug(f"Encoding {encoding} failed for {filename}: {str(e)}")
                last_error = e
                continue
            except Exception as e:
                logger.warning(f"⚠️  Error with encoding {encoding} for {filename}: {e}")
                last_error = e
                continue
        
        # If we get here, no encoding worked
        logger.error(f"❌ Failed to decode {filename} with any encoding. Attempted: {', '.join(encodings_to_try)}")
        
        raise DocumentProcessingError(
            ErrorMessageTemplates.encoding_error(filename, encodings_to_try),
            error_category=ErrorCategory.ENCODING_ERROR,
            is_user_error=True,
            details={
                "filename": filename,
                "file_type": "text",
                "attempted_encodings": encodings_to_try,
                "last_error": str(last_error) if last_error else None
            }
        )
    
    async def _load_docx(
        self, 
        file_path: str, 
        filename: str,
        image_analyzer: Optional[Any] = None
    ) -> List[Document]:
        """Load DOCX document using Docx2txtLoader with optional image analysis"""
        try:
            loader = Docx2txtLoader(file_path)
            documents = loader.load()
            
            if not documents or not documents[0].page_content.strip():
                raise DocumentProcessingError(
                    ErrorMessageTemplates.empty_content(filename),
                    error_category=ErrorCategory.EMPTY_CONTENT,
                    is_user_error=True,
                    details={"filename": filename, "file_type": "docx"}
                )
            
            # Clean and enhance the document
            doc = documents[0]
            doc.page_content = self._clean_text(doc.page_content)
            doc.metadata.update({
                "source": filename,
                "loader": "Docx2txtLoader"
            })
            
            # --- MULTIMODAL IMAGE EXTRACTION ---
            if image_analyzer:
                try:
                    import docx
                    
                    doc_obj = docx.Document(file_path)
                    image_descriptions = []
                    image_count = 0
                    
                    # Extract images from relationships
                    # This finds all images in the document
                    for rel in doc_obj.part.rels.values():
                        if "image" in rel.target_ref:
                            try:
                                img_part = rel.target_part
                                img_bytes = img_part.blob
                                
                                # Analyze image
                                description = await image_analyzer(img_bytes)
                                if description:
                                    image_count += 1
                                    image_descriptions.append(f"[IMAGE {image_count}: {description}]")
                            except Exception as img_err:
                                logger.warning(f"Failed to analyze image in DOCX: {img_err}")
                    
                    if image_descriptions:
                        logger.info(f"🖼️ Found {len(image_descriptions)} images in DOCX: {filename}")
                        # Append image descriptions to content
                        doc.page_content += "\n\n=== IMAGES IN DOCUMENT ===\n" + "\n\n".join(image_descriptions)
                        doc.metadata["has_images"] = True
                        doc.metadata["image_count"] = image_count
                        
                except Exception as e:
                    logger.warning(f"Multimodal extraction failed for DOCX {filename}: {e}")
            
            logger.debug(f"Successfully loaded DOCX file: {filename}")
            return [doc]
            
        except DocumentProcessingError:
            raise
        except Exception as e:
            logger.error(f"❌ DOCX processing error for {filename}: {e}", exc_info=True)
            if "docx" in str(e).lower() or "corrupt" in str(e).lower() or "invalid" in str(e).lower():
                raise DocumentProcessingError(
                    ErrorMessageTemplates.corrupted_file(filename, "DOCX"),
                    error_category=ErrorCategory.CORRUPTED_FILE,
                    is_user_error=True,
                    details={"filename": filename, "file_type": "docx", "error": str(e)}
                )
            raise DocumentProcessingError(
                ErrorMessageTemplates.processing_error(filename, str(e)),
                error_category=ErrorCategory.PROCESSING_ERROR,
                is_user_error=False,
                details={"filename": filename, "file_type": "docx", "error": str(e)}
            )
    
    async def _load_pptx(
        self, 
        file_path: str, 
        filename: str,
        image_analyzer: Optional[Any] = None
    ) -> List[Document]:
        """Load PPTX document using UnstructuredPowerPointLoader with optional image analysis"""
        try:
            loader = UnstructuredPowerPointLoader(file_path)
            documents = loader.load()
            
            if not documents or not documents[0].page_content.strip():
                raise DocumentProcessingError(
                    ErrorMessageTemplates.empty_content(filename),
                    error_category=ErrorCategory.EMPTY_CONTENT,
                    is_user_error=True,
                    details={"filename": filename, "file_type": "pptx"}
                )
            
            # Clean and enhance the document
            doc = documents[0]
            doc.page_content = self._clean_text(doc.page_content)
            doc.metadata.update({
                "source": filename,
                "loader": "UnstructuredPowerPointLoader"
            })
            
            # --- MULTIMODAL IMAGE EXTRACTION ---
            if image_analyzer:
                try:
                    from pptx import Presentation
                    from pptx.enum.shapes import MSO_SHAPE_TYPE
                    
                    prs = Presentation(file_path)
                    image_descriptions = []
                    total_images = 0
                    
                    for i, slide in enumerate(prs.slides):
                        slide_images = []
                        for shape in slide.shapes:
                            try:
                                # Check if shape is a picture
                                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                    image = shape.image
                                    img_bytes = image.blob
                                    
                                    # Analyze image
                                    description = await image_analyzer(img_bytes)
                                    if description:
                                        slide_images.append(description)
                                        total_images += 1
                            except Exception as shape_err:
                                continue
                        
                        if slide_images:
                            # Append slide images to content, marking the slide number
                            # Note: Unstructured loader usually dumps all text. We'll append images at the end.
                            image_descriptions.append(f"\n--- Images on Slide {i+1} ---")
                            for desc in slide_images:
                                image_descriptions.append(f"[IMAGE: {desc}]")
                    
                    if image_descriptions:
                        logger.info(f"🖼️ Found {total_images} images in PPTX: {filename}")
                        doc.page_content += "\n\n=== SLIDE IMAGES ===\n" + "\n".join(image_descriptions)
                        doc.metadata["has_images"] = True
                        doc.metadata["image_count"] = total_images
                        
                except Exception as e:
                    logger.warning(f"Multimodal extraction failed for PPTX {filename}: {e}")
            
            logger.debug(f"Successfully loaded PPTX file: {filename}")
            return [doc]
            
        except DocumentProcessingError:
            raise
        except Exception as e:
            logger.error(f"❌ PPTX processing error for {filename}: {e}", exc_info=True)
            if "pptx" in str(e).lower() or "corrupt" in str(e).lower() or "invalid" in str(e).lower():
                raise DocumentProcessingError(
                    ErrorMessageTemplates.corrupted_file(filename, "PPTX"),
                    error_category=ErrorCategory.CORRUPTED_FILE,
                    is_user_error=True,
                    details={"filename": filename, "file_type": "pptx", "error": str(e)}
                )
            raise DocumentProcessingError(
                ErrorMessageTemplates.processing_error(filename, str(e)),
                error_category=ErrorCategory.PROCESSING_ERROR,
                is_user_error=False,
                details={"filename": filename, "file_type": "pptx", "error": str(e)}
            )
    
    async def _load_xlsx(self, file_path: str, filename: str) -> List[Document]:
        """Load XLSX document using pandas"""
        try:
            # Read all sheets from the Excel file
            excel_file = pd.ExcelFile(file_path, engine='openpyxl')
            documents = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # Convert dataframe to text representation
                # Include column headers and data
                content_parts = [f"Sheet: {sheet_name}\n"]
                
                # Add column headers
                content_parts.append("Columns: " + ", ".join(df.columns.astype(str)) + "\n\n")
                
                # Add data rows
                for idx, row in df.iterrows():
                    row_text = " | ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                    if row_text:
                        content_parts.append(row_text + "\n")
                
                content = "".join(content_parts)
                
                if content.strip():
                    doc = Document(
                        page_content=self._clean_text(content),
                        metadata={
                            "source": filename,
                            "sheet_name": sheet_name,
                            "loader": "PandasExcelLoader",
                            "rows": len(df),
                            "columns": len(df.columns)
                        }
                    )
                    documents.append(doc)
            
            if not documents:
                raise DocumentProcessingError(
                    ErrorMessageTemplates.empty_content(filename),
                    error_category=ErrorCategory.EMPTY_CONTENT,
                    is_user_error=True,
                    details={"filename": filename, "file_type": "xlsx"}
                )
            
            logger.debug(f"Successfully loaded XLSX file: {filename} ({len(documents)} sheets)")
            return documents
            
        except DocumentProcessingError:
            raise
        except Exception as e:
            logger.error(f"❌ XLSX processing error for {filename}: {e}", exc_info=True)
            if "xlsx" in str(e).lower() or "corrupt" in str(e).lower() or "excel" in str(e).lower():
                raise DocumentProcessingError(
                    ErrorMessageTemplates.corrupted_file(filename, "XLSX"),
                    error_category=ErrorCategory.CORRUPTED_FILE,
                    is_user_error=True,
                    details={"filename": filename, "file_type": "xlsx", "error": str(e)}
                )
            raise DocumentProcessingError(
                ErrorMessageTemplates.processing_error(filename, str(e)),
                error_category=ErrorCategory.PROCESSING_ERROR,
                is_user_error=False,
                details={"filename": filename, "file_type": "xlsx", "error": str(e)}
            )
    
    async def _load_csv(self, file_path: str, filename: str) -> List[Document]:
        """Load CSV document using pandas"""
        encodings_to_try = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        try:
            # Try different encodings for CSV
            df = None
            successful_encoding = None
            
            for encoding in encodings_to_try:
                try:
                    df = pd.read_csv(file_path, encoding=encoding)
                    successful_encoding = encoding
                    logger.debug(f"Successfully decoded CSV {filename} with encoding: {encoding}")
                    break
                except UnicodeDecodeError:
                    continue
            
            if df is None:
                raise DocumentProcessingError(
                    ErrorMessageTemplates.encoding_error(filename, encodings_to_try),
                    error_category=ErrorCategory.ENCODING_ERROR,
                    is_user_error=True,
                    details={
                        "filename": filename,
                        "file_type": "csv",
                        "attempted_encodings": encodings_to_try
                    }
                )
            
            # Convert dataframe to text representation
            content_parts = []
            
            # Add column headers
            content_parts.append("Columns: " + ", ".join(df.columns.astype(str)) + "\n\n")
            
            # Add data rows
            for idx, row in df.iterrows():
                row_text = " | ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                if row_text:
                    content_parts.append(row_text + "\n")
            
            content = "".join(content_parts)
            
            if not content.strip():
                raise DocumentProcessingError(
                    ErrorMessageTemplates.empty_content(filename),
                    error_category=ErrorCategory.EMPTY_CONTENT,
                    is_user_error=True,
                    details={"filename": filename, "file_type": "csv"}
                )
            
            doc = Document(
                page_content=self._clean_text(content),
                metadata={
                    "source": filename,
                    "loader": "PandasCSVLoader",
                    "rows": len(df),
                    "columns": len(df.columns),
                    "encoding": successful_encoding
                }
            )
            
            logger.debug(f"Successfully loaded CSV file: {filename} ({len(df)} rows, {len(df.columns)} columns)")
            return [doc]
            
        except DocumentProcessingError:
            raise
        except Exception as e:
            logger.error(f"❌ CSV processing error for {filename}: {e}", exc_info=True)
            if "csv" in str(e).lower() or "corrupt" in str(e).lower() or "parse" in str(e).lower():
                raise DocumentProcessingError(
                    ErrorMessageTemplates.corrupted_file(filename, "CSV"),
                    error_category=ErrorCategory.CORRUPTED_FILE,
                    is_user_error=True,
                    details={"filename": filename, "file_type": "csv", "error": str(e)}
                )
            raise DocumentProcessingError(
                ErrorMessageTemplates.processing_error(filename, str(e)),
                error_category=ErrorCategory.PROCESSING_ERROR,
                is_user_error=False,
                details={"filename": filename, "file_type": "csv", "error": str(e)}
            )
    
    async def _load_image(self, file_path: str, filename: str) -> List[Document]:
        """
        Load image document - creates a document with image path for embedding
        The actual image embedding will be handled by Cohere's multimodal embeddings
        """
        try:
            # Verify image file exists and is readable
            if not os.path.exists(file_path):
                raise DocumentProcessingError(
                    ErrorMessageTemplates.file_not_found(filename),
                    error_category=ErrorCategory.FILE_NOT_FOUND,
                    is_user_error=False,
                    details={"filename": filename, "file_type": "image"}
                )
            
            # Get image file size
            file_size = os.path.getsize(file_path)
            if file_size == 0:
                raise DocumentProcessingError(
                    ErrorMessageTemplates.empty_content(filename),
                    error_category=ErrorCategory.EMPTY_CONTENT,
                    is_user_error=True,
                    details={"filename": filename, "file_type": "image", "file_size": 0}
                )
            
            # Create a document with image metadata
            # The content will be a placeholder that indicates this is an image
            doc = Document(
                page_content=f"[IMAGE: {filename}]",
                metadata={
                    "source": filename,
                    "loader": "ImageLoader",
                    "file_type": "image",
                    "image_path": file_path,
                    "file_size": file_size,
                    "extension": Path(filename).suffix.lower()
                }
            )
            
            logger.debug(f"Successfully loaded image file: {filename} ({file_size} bytes)")
            return [doc]
            
        except DocumentProcessingError:
            raise
        except Exception as e:
            logger.error(f"❌ Image processing error for {filename}: {e}", exc_info=True)
            raise DocumentProcessingError(
                ErrorMessageTemplates.processing_error(filename, str(e)),
                error_category=ErrorCategory.PROCESSING_ERROR,
                is_user_error=False,
                details={"filename": filename, "file_type": "image", "error": str(e)}
            )
    
    async def _load_sdf(self, file_path: str, filename: str) -> List[Document]:
        """
        Load SDF (Structure Data File) document
        
        Supports chemical structure files commonly used in pharmaceutical research.
        Attempts to use rdkit for advanced parsing, falls back to basic parsing if unavailable.
        
        Args:
            file_path: Path to temporary SDF file
            filename: Original filename
            
        Returns:
            List of Document objects, one per molecular structure
            
        Raises:
            DocumentProcessingError: If SDF parsing fails
        """
        try:
            # Try to import rdkit for advanced parsing
            try:
                from rdkit import Chem
                from rdkit.Chem import Descriptors
                use_rdkit = True
                logger.debug("Using rdkit for SDF parsing")
            except ImportError:
                use_rdkit = False
                logger.debug("rdkit not available, using fallback SDF parser")
            
            documents = []
            
            if use_rdkit:
                # Use rdkit for robust parsing
                try:
                    supplier = Chem.SDMolSupplier(file_path, removeHs=False, sanitize=True)
                except Exception as e:
                    logger.error(f"❌ Failed to create SDF supplier for {filename}: {e}", exc_info=True)
                    raise DocumentProcessingError(
                        ErrorMessageTemplates.corrupted_file(filename, "SDF"),
                        error_category=ErrorCategory.CORRUPTED_FILE,
                        is_user_error=True,
                        details={"filename": filename, "file_type": "sdf", "error": str(e)}
                    )
                
                for idx, mol in enumerate(supplier):
                    if mol is None:
                        logger.warning(f"⚠️  Skipping invalid molecule at index {idx} in {filename}")
                        continue
                    
                    try:
                        # Extract molecular properties
                        properties = mol.GetPropsAsDict()
                        
                        # Get basic molecular information
                        mol_formula = Chem.rdMolDescriptors.CalcMolFormula(mol)
                        mol_weight = Descriptors.MolWt(mol)
                        num_atoms = mol.GetNumAtoms()
                        num_bonds = mol.GetNumBonds()
                        
                        # Try to get compound name from properties or use index
                        compound_name = properties.get('_Name', properties.get('Name', f"Compound_{idx + 1}"))
                        
                        # Build text representation
                        content_parts = [
                            f"Compound: {compound_name}",
                            f"Molecular Formula: {mol_formula}",
                            f"Molecular Weight: {mol_weight:.2f} g/mol",
                            f"Number of Atoms: {num_atoms}",
                            f"Number of Bonds: {num_bonds}",
                        ]
                        
                        # Add SMILES representation if possible
                        try:
                            smiles = Chem.MolToSmiles(mol)
                            content_parts.append(f"SMILES: {smiles}")
                        except Exception as e:
                            logger.debug(f"Could not generate SMILES for molecule {idx}: {e}")
                        
                        # Add InChI if possible
                        try:
                            inchi = Chem.MolToInchi(mol)
                            content_parts.append(f"InChI: {inchi}")
                        except Exception as e:
                            logger.debug(f"Could not generate InChI for molecule {idx}: {e}")
                        
                        # Add additional properties
                        if properties:
                            content_parts.append("\nProperties:")
                            for key, value in properties.items():
                                if not key.startswith('_'):  # Skip internal properties
                                    content_parts.append(f"  {key}: {value}")
                        
                        page_content = "\n".join(content_parts)
                        
                        # Create document with metadata
                        doc = Document(
                            page_content=page_content,
                            metadata={
                                "source": filename,
                                "loader": "SDFLoader",
                                "file_type": "sdf",
                                "compound_name": compound_name,
                                "molecular_formula": mol_formula,
                                "molecular_weight": mol_weight,
                                "num_atoms": num_atoms,
                                "num_bonds": num_bonds,
                                "structure_index": idx,
                                "properties": properties
                            }
                        )
                        documents.append(doc)
                        
                    except Exception as e:
                        logger.warning(f"⚠️  Error processing molecule {idx} in {filename}: {e}", exc_info=True)
                        continue
                
            else:
                # Fallback: Basic text parsing without rdkit
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                
                # Split by molecule delimiter ($$$$)
                molecules = content.split('$$$$')
                
                for idx, mol_block in enumerate(molecules):
                    mol_block = mol_block.strip()
                    if not mol_block:
                        continue
                    
                    # Parse basic information from SDF format
                    lines = mol_block.split('\n')
                    
                    # First line is typically the compound name
                    compound_name = lines[0].strip() if lines else f"Compound_{idx + 1}"
                    if not compound_name:
                        compound_name = f"Compound_{idx + 1}"
                    
                    # Extract properties from data fields (lines starting with >)
                    properties = {}
                    content_parts = [f"Compound: {compound_name}"]
                    
                    i = 0
                    while i < len(lines):
                        line = lines[i].strip()
                        if line.startswith('>'):
                            # Property field
                            prop_name = line.strip('<> ')
                            if i + 1 < len(lines):
                                prop_value = lines[i + 1].strip()
                                properties[prop_name] = prop_value
                                content_parts.append(f"{prop_name}: {prop_value}")
                                i += 2
                            else:
                                i += 1
                        else:
                            i += 1
                    
                    # Try to extract atom and bond counts from counts line (line 4 in SDF format)
                    if len(lines) > 3:
                        counts_line = lines[3].strip()
                        parts = counts_line.split()
                        if len(parts) >= 2:
                            try:
                                num_atoms = int(parts[0])
                                num_bonds = int(parts[1])
                                content_parts.insert(1, f"Number of Atoms: {num_atoms}")
                                content_parts.insert(2, f"Number of Bonds: {num_bonds}")
                            except (ValueError, IndexError):
                                num_atoms = None
                                num_bonds = None
                        else:
                            num_atoms = None
                            num_bonds = None
                    else:
                        num_atoms = None
                        num_bonds = None
                    
                    if properties:
                        content_parts.append("\nAdditional Properties:")
                        for key, value in properties.items():
                            content_parts.append(f"  {key}: {value}")
                    
                    page_content = "\n".join(content_parts)
                    
                    # Create document with available metadata
                    metadata = {
                        "source": filename,
                        "loader": "SDFLoader",
                        "file_type": "sdf",
                        "compound_name": compound_name,
                        "structure_index": idx,
                        "properties": properties,
                        "parsing_method": "fallback"
                    }
                    
                    if num_atoms is not None:
                        metadata["num_atoms"] = num_atoms
                    if num_bonds is not None:
                        metadata["num_bonds"] = num_bonds
                    
                    doc = Document(
                        page_content=page_content,
                        metadata=metadata
                    )
                    documents.append(doc)
            
            if not documents:
                raise DocumentProcessingError(
                    ErrorMessageTemplates.empty_content(filename),
                    error_category=ErrorCategory.EMPTY_CONTENT,
                    is_user_error=True,
                    details={
                        "filename": filename,
                        "file_type": "sdf",
                        "reason": "No valid molecular structures found"
                    }
                )
            
            logger.info(f"✅ Successfully loaded {len(documents)} molecular structure(s) from {filename}")
            return documents
            
        except DocumentProcessingError:
            raise
        except Exception as e:
            logger.error(f"❌ Error loading SDF file {filename}: {e}", exc_info=True)
            # Check if it's a corrupted file error
            if "corrupt" in str(e).lower() or "invalid" in str(e).lower() or "parse" in str(e).lower():
                raise DocumentProcessingError(
                    ErrorMessageTemplates.corrupted_file(filename, "SDF"),
                    error_category=ErrorCategory.CORRUPTED_FILE,
                    is_user_error=True,
                    details={"filename": filename, "file_type": "sdf", "error": str(e)}
                )
            raise DocumentProcessingError(
                ErrorMessageTemplates.processing_error(filename, str(e)),
                error_category=ErrorCategory.PROCESSING_ERROR,
                is_user_error=False,
                details={"filename": filename, "file_type": "sdf", "error": str(e)}
            )
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize extracted text"""
        if not text:
            return ""
        
        # Remove excessive whitespace
        import re
        
        # Replace multiple spaces with single space
        text = re.sub(r' +', ' ', text)
        
        # Replace multiple newlines with double newline
        text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
        
        # Remove trailing whitespace from lines
        lines = [line.rstrip() for line in text.split('\n')]
        text = '\n'.join(lines)
        
        # Remove leading/trailing whitespace
        text = text.strip()
        
        return text
    
    async def validate_document_content(
        self, 
        documents: List[Document], 
        filename: str = "document",
        file_type: str = "unknown"
    ) -> Dict[str, Any]:
        """
        Validate loaded document content and return structured validation results
        
        Enhanced validation with specific error messages for different content types
        and failure reasons. Minimum content threshold lowered to 10 characters.
        
        Args:
            documents: List of Document objects to validate
            filename: Original filename for error messages
            file_type: File type/extension for specific validation rules
            
        Returns:
            Dictionary with validation results including:
            - valid: Boolean indicating if content passes validation
            - error: Specific error message if validation fails (only if valid=False)
            - failure_reason: Categorized failure reason (only if valid=False)
            - warnings: List of non-critical issues
            - stats: Content statistics
            - content_type: Detected content type for type-specific handling
        """
        start_time = time.time()
        
        # Log validation start
        logger.debug(
            f"🔍 Starting content validation for {filename} ({len(documents) if documents else 0} documents)",
            extra={
                'operation': 'content_validation',
                'document_name': filename,
                'file_type': file_type,
                'document_count': len(documents) if documents else 0
            }
        )
        
        # Minimum content threshold (lowered to 10 characters per requirements)
        MIN_CONTENT_THRESHOLD = 10
        
        # Initialize validation result structure
        validation_result = {
            "valid": True,
            "warnings": [],
            "stats": {},
            "content_type": self._detect_content_type(file_type)
        }
        
        # Validation Rule 0: Check if documents list is empty
        if not documents:
            validation_result["valid"] = False
            validation_result["error"] = (
                f"No documents were extracted from '{filename}'. "
                f"The file may be corrupted, empty, or in an unsupported format."
            )
            validation_result["failure_reason"] = "no_documents"
            validation_result["stats"] = {
                "document_count": 0,
                "total_characters": 0,
                "total_words": 0,
                "pages_with_content": 0
            }
            duration = time.time() - start_time
            logger.warning(
                f"⚠️  Validation failed for {filename}: No documents extracted",
                extra={
                    'operation': 'content_validation',
                    'document_name': filename,
                    'file_type': file_type,
                    'validation_result': 'failed',
                    'failure_reason': 'no_documents',
                    'duration': duration * 1000
                }
            )
            return validation_result
        
        # Collect comprehensive content statistics
        stats = {
            "document_count": len(documents),
            "total_characters": 0,
            "total_words": 0,
            "average_page_length": 0,
            "empty_pages": 0,
            "pages_with_content": 0,
            "non_whitespace_chars": 0,
            "min_page_length": float('inf'),
            "max_page_length": 0,
            "content_type": validation_result["content_type"]
        }
        
        # Collect content statistics from all documents
        for doc in documents:
            content = doc.page_content.strip()
            content_length = len(content)
            
            if content:
                stats["pages_with_content"] += 1
                stats["total_characters"] += content_length
                stats["total_words"] += len(content.split())
                stats["non_whitespace_chars"] += len(''.join(content.split()))
                stats["min_page_length"] = min(stats["min_page_length"], content_length)
                stats["max_page_length"] = max(stats["max_page_length"], content_length)
            else:
                stats["empty_pages"] += 1
        
        # Calculate average page length
        if stats["pages_with_content"] > 0:
            stats["average_page_length"] = stats["total_characters"] / stats["pages_with_content"]
        else:
            stats["min_page_length"] = 0
        
        validation_result["stats"] = stats
        
        # Validation Rule 1: Check for completely empty content
        if stats["pages_with_content"] == 0:
            validation_result["valid"] = False
            validation_result["error"] = (
                f"No content could be extracted from '{filename}'. "
                f"The file appears to be empty or contains only non-text elements. "
                f"For {validation_result['content_type']} files, ensure the file contains readable text or data."
            )
            validation_result["failure_reason"] = "empty_content"
            duration = time.time() - start_time
            logger.warning(
                f"⚠️  Validation failed for {filename}: Empty content "
                f"(document_count={stats['document_count']}, pages_with_content=0)",
                extra={
                    'operation': 'content_validation',
                    'document_name': filename,
                    'file_type': file_type,
                    'validation_result': 'failed',
                    'failure_reason': 'empty_content',
                    'document_count': stats['document_count'],
                    'pages_with_content': 0,
                    'duration': duration * 1000
                }
            )
            return validation_result
        
        # Validation Rule 2: Check for whitespace-only content
        if stats["non_whitespace_chars"] == 0:
            validation_result["valid"] = False
            validation_result["error"] = (
                f"No readable content found in '{filename}'. "
                f"The file contains only whitespace characters. "
                f"Please ensure the {validation_result['content_type']} file has actual text or data content."
            )
            validation_result["failure_reason"] = "whitespace_only"
            duration = time.time() - start_time
            logger.warning(
                f"⚠️  Validation failed for {filename}: Whitespace-only content "
                f"(total_chars={stats['total_characters']}, non_whitespace=0)",
                extra={
                    'operation': 'content_validation',
                    'document_name': filename,
                    'file_type': file_type,
                    'validation_result': 'failed',
                    'failure_reason': 'whitespace_only',
                    'total_characters': stats['total_characters'],
                    'non_whitespace_chars': 0,
                    'duration': duration * 1000
                }
            )
            return validation_result
        
        # Validation Rule 3: Check minimum content threshold
        if stats["total_characters"] < MIN_CONTENT_THRESHOLD:
            validation_result["valid"] = False
            validation_result["error"] = (
                f"Insufficient content in '{filename}'. "
                f"Found {stats['total_characters']} characters, minimum required is {MIN_CONTENT_THRESHOLD}. "
                f"Please provide a {validation_result['content_type']} file with more substantial content."
            )
            validation_result["failure_reason"] = "insufficient_content"
            validation_result["stats"]["min_required"] = MIN_CONTENT_THRESHOLD
            duration = time.time() - start_time
            logger.warning(
                f"⚠️  Validation failed for {filename}: Insufficient content "
                f"(chars={stats['total_characters']}, min={MIN_CONTENT_THRESHOLD})",
                extra={
                    'operation': 'content_validation',
                    'document_name': filename,
                    'file_type': file_type,
                    'validation_result': 'failed',
                    'failure_reason': 'insufficient_content',
                    'total_characters': stats['total_characters'],
                    'min_required': MIN_CONTENT_THRESHOLD,
                    'duration': duration * 1000
                }
            )
            return validation_result
        
        # Content passed basic validation - now check for warnings
        
        # Warning 1: Very little content (between 10-50 characters)
        if stats["total_characters"] < 50:
            validation_result["warnings"].append(
                f"'{filename}' contains very little text content ({stats['total_characters']} characters). "
                f"The document was processed but may not provide useful context for queries."
            )
            logger.info(f"ℹ️  Low content warning for {filename}: {stats['total_characters']} characters")
        
        # Warning 2: Many empty pages
        if stats["empty_pages"] > 0 and stats["empty_pages"] > stats["pages_with_content"]:
            validation_result["warnings"].append(
                f"More than half of the pages in '{filename}' are empty "
                f"({stats['empty_pages']} empty, {stats['pages_with_content']} with content). "
                f"Consider reviewing the source file."
            )
        
        # Warning 3: Very short average page length
        if stats["average_page_length"] < 20 and stats["pages_with_content"] > 1:
            validation_result["warnings"].append(
                f"Pages in '{filename}' have very short content "
                f"(average {stats['average_page_length']:.0f} characters per page). "
                f"This may indicate extraction issues."
            )
        
        # Type-specific validation for different content types
        self._apply_content_type_validation(validation_result, documents, filename, file_type, stats)
        
        # Log successful validation
        duration = time.time() - start_time
        if validation_result["valid"]:
            logger.info(
                f"✅ Validation passed for {filename}: "
                f"{stats['total_characters']} chars, {stats['total_words']} words, "
                f"{stats['pages_with_content']} pages ({duration:.3f}s)",
                extra={
                    'operation': 'content_validation',
                    'document_name': filename,
                    'file_type': file_type,
                    'validation_result': 'passed',
                    'total_characters': stats['total_characters'],
                    'total_words': stats['total_words'],
                    'pages_with_content': stats['pages_with_content'],
                    'warning_count': len(validation_result['warnings']),
                    'duration': duration * 1000
                }
            )
        
        return validation_result
    
    def _detect_content_type(self, file_type: str) -> str:
        """
        Detect and categorize content type for type-specific validation
        
        Args:
            file_type: File extension (e.g., '.pdf', '.csv')
            
        Returns:
            Human-readable content type category
        """
        content_type_map = {
            '.pdf': 'PDF document',
            '.docx': 'Word document',
            '.txt': 'text file',
            '.md': 'Markdown document',
            '.pptx': 'PowerPoint presentation',
            '.xlsx': 'Excel spreadsheet',
            '.csv': 'CSV data file',
            '.sdf': 'chemical structure file',
            '.mol': 'molecular structure file',
            '.png': 'image file',
            '.jpg': 'image file',
            '.jpeg': 'image file',
            '.gif': 'image file',
            '.bmp': 'image file',
            '.webp': 'image file'
        }
        return content_type_map.get(file_type.lower(), 'document')
    
    def _apply_content_type_validation(
        self,
        validation_result: Dict[str, Any],
        documents: List[Document],
        filename: str,
        file_type: str,
        stats: Dict[str, Any]
    ) -> None:
        """
        Apply content-type-specific validation rules and add warnings
        
        Args:
            validation_result: Validation result dictionary to update
            documents: List of documents being validated
            filename: Original filename
            file_type: File extension
            stats: Content statistics
        """
        # Document types (PDF, DOCX, PPTX) - expect substantial text content
        if file_type in ['.pdf', '.docx', '.pptx']:
            if stats["total_words"] < 5:
                validation_result["warnings"].append(
                    f"'{filename}' contains fewer than 5 words ({stats['total_words']} words). "
                    f"This may indicate a scanning issue, image-only content, or extraction problem. "
                    f"For scanned documents, consider using OCR."
                )
            
            # Check for very short documents
            if stats["total_words"] < 20 and stats["total_characters"] >= 10:
                validation_result["warnings"].append(
                    f"'{filename}' has minimal text content ({stats['total_words']} words). "
                    f"Verify this is the complete document."
                )
        
        # Data files (CSV, XLSX) - check for structured content
        elif file_type in ['.csv', '.xlsx']:
            if stats["total_words"] < 3:
                validation_result["warnings"].append(
                    f"'{filename}' appears to have minimal data ({stats['total_words']} words). "
                    f"Verify the file contains the expected tabular information. "
                    f"Empty or header-only files may not be useful for queries."
                )
            
            # Check if it looks like headers only
            if stats["pages_with_content"] == 1 and stats["total_words"] < 10:
                validation_result["warnings"].append(
                    f"'{filename}' may contain only column headers with no data rows. "
                    f"Ensure the file includes actual data."
                )
        
        # Chemical structure files (SDF, MOL) - validate molecular data
        elif file_type in ['.sdf', '.mol']:
            has_molecular_data = any(
                'Molecular Formula' in doc.page_content or 
                'Compound' in doc.page_content or
                'SMILES' in doc.page_content
                for doc in documents
            )
            
            if not has_molecular_data:
                validation_result["warnings"].append(
                    f"'{filename}' may not contain valid molecular structure data. "
                    f"Expected fields like 'Molecular Formula' or 'Compound' were not found. "
                    f"Verify the file is a valid SDF/MOL format."
                )
            
            # Check for multiple structures
            if len(documents) > 1:
                validation_result["warnings"].append(
                    f"'{filename}' contains {len(documents)} molecular structures. "
                    f"Each structure will be processed separately."
                )
        
        # Text files (TXT, MD) - basic content checks
        elif file_type in ['.txt', '.md']:
            if stats["total_words"] < 10 and stats["total_characters"] >= 10:
                validation_result["warnings"].append(
                    f"'{filename}' has very few words ({stats['total_words']} words) "
                    f"despite having {stats['total_characters']} characters. "
                    f"The file may contain special characters or non-standard formatting."
                )
        
        # Image files - special handling
        elif file_type in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp']:
            # Images are handled differently - they should have placeholder content
            if not any('[IMAGE:' in doc.page_content for doc in documents):
                validation_result["warnings"].append(
                    f"'{filename}' is an image file but may not have been processed correctly. "
                    f"Image content should be handled by multimodal embeddings."
                )
        
        # Log type-specific validation results
        if validation_result["warnings"]:
            logger.debug(
                f"Type-specific validation for {filename} ({file_type}): "
                f"{len(validation_result['warnings'])} warnings"
            )
    
    def get_processing_stats(self) -> Dict[str, Any]:
        """Get processing statistics and capabilities"""
        return {
            "supported_formats": self.get_supported_formats(),
            "langchain_version": "0.1.0",
            "loaders": {
                ".pdf": "PyPDFLoader",
                ".txt": "TextLoader",
                ".md": "TextLoader", 
                ".docx": "Docx2txtLoader",
                ".pptx": "UnstructuredPowerPointLoader",
                ".xlsx": "PandasExcelLoader",
                ".csv": "PandasCSVLoader",
                ".sdf": "SDFLoader",
                ".mol": "SDFLoader"
            },
            "features": [
                "Encoding detection for text files",
                "Page-by-page PDF processing",
                "Multi-sheet Excel processing",
                "CSV with encoding detection",
                "PowerPoint slide extraction",
                "SDF/MOL chemical structure parsing",
                "Multi-molecule SDF support",
                "Text cleaning and normalization",
                "Metadata enhancement",
                "Content validation"
            ]
        }


# Global document loader instance
document_loader = EnhancedDocumentLoader()