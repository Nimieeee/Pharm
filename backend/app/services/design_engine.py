"""
Design Intelligence Engine

Rule-based design system for professional slide generation.
Contains 10 curated color themes, typography constants, and layout rules.

The AI generates CONTENT. This engine handles VISUAL DESIGN.
Never let the AI choose fonts, sizes, or colors directly.
"""

import io
from typing import Dict, List, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData


class DesignEngine:
    """
    Rule-based design system for slide generation.
    
    Features:
    - 10 curated color themes with professional palettes
    - Typography hierarchy (44pt/28pt/16pt/11pt)
    - Layout engine with anti-repetition rules
    - Density analysis (max 6 bullets per slide)
    - PPTX assembly with professional formatting
    """
    
    # ========================================
    # THEME DEFINITIONS (10 curated palettes)
    # ========================================
    
    THEMES = {
        "ocean_gradient": {
            "name": "Ocean Gradient",
            "bg_dark": "065A82",      # Title slide background
            "bg_light": "F8FBFD",     # Content slide background
            "primary": "1C7293",      # Section headers, accents
            "accent": "21295C",       # Callout backgrounds, icons
            "text_on_dark": "FFFFFF", # Text on dark backgrounds
            "text_on_light": "2D3436",# Body text on light backgrounds
            "muted": "636E72",        # Captions, metadata
            "success": "00B894",      # Positive data callouts
            "warning": "FDCB6E",      # Warning highlights
            "header_font": "Georgia",
            "body_font": "Calibri",
        },
        "forest_moss": {
            "name": "Forest & Moss",
            "bg_dark": "2C5F2D",
            "bg_light": "F5F5F0",
            "primary": "97BC62",
            "accent": "2C5F2D",
            "text_on_dark": "FFFFFF",
            "text_on_light": "2D3436",
            "muted": "636E72",
            "success": "00B894",
            "warning": "FDCB6E",
            "header_font": "Cambria",
            "body_font": "Calibri",
        },
        "coral_energy": {
            "name": "Coral Energy",
            "bg_dark": "2F3C7E",
            "bg_light": "FFF9F5",
            "primary": "F96167",
            "accent": "F9E795",
            "text_on_dark": "FFFFFF",
            "text_on_light": "2D3436",
            "muted": "636E72",
            "success": "00B894",
            "warning": "FDCB6E",
            "header_font": "Arial Black",
            "body_font": "Arial",
        },
        "teal_trust": {
            "name": "Teal Trust",
            "bg_dark": "028090",
            "bg_light": "F0FAFA",
            "primary": "00A896",
            "accent": "02C39A",
            "text_on_dark": "FFFFFF",
            "text_on_light": "2D3436",
            "muted": "636E72",
            "success": "02C39A",
            "warning": "FDCB6E",
            "header_font": "Trebuchet MS",
            "body_font": "Calibri",
        },
        "charcoal_minimal": {
            "name": "Charcoal Minimal",
            "bg_dark": "212121",
            "bg_light": "F8F8F8",
            "primary": "36454F",
            "accent": "F2F2F2",
            "text_on_dark": "F2F2F2",
            "text_on_light": "212121",
            "muted": "888888",
            "success": "4CAF50",
            "warning": "FF9800",
            "header_font": "Arial Black",
            "body_font": "Calibri Light",
        },
        "berry_cream": {
            "name": "Berry Cream",
            "bg_dark": "6B2737",
            "bg_light": "FFF5F7",
            "primary": "D63384",
            "accent": "F8D7DA",
            "text_on_dark": "FFFFFF",
            "text_on_light": "2D3436",
            "muted": "636E72",
            "success": "00B894",
            "warning": "FDCB6E",
            "header_font": "Georgia",
            "body_font": "Calibri",
        },
        "sage_calm": {
            "name": "Sage Calm",
            "bg_dark": "556B2F",
            "bg_light": "F5F5F0",
            "primary": "8FBC8F",
            "accent": "DEDEB8",
            "text_on_dark": "FFFFFF",
            "text_on_light": "2D3436",
            "muted": "636E72",
            "success": "00B894",
            "warning": "FDCB6E",
            "header_font": "Cambria",
            "body_font": "Calibri",
        },
        "cherry_bold": {
            "name": "Cherry Bold",
            "bg_dark": "8B0000",
            "bg_light": "FFFAFA",
            "primary": "DC143C",
            "accent": "FFB6C1",
            "text_on_dark": "FFFFFF",
            "text_on_light": "2D3436",
            "muted": "636E72",
            "success": "00B894",
            "warning": "FDCB6E",
            "header_font": "Arial Black",
            "body_font": "Arial",
        },
        "midnight_executive": {
            "name": "Midnight Executive",
            "bg_dark": "1A1A2E",
            "bg_light": "F8F9FA",
            "primary": "16213E",
            "accent": "0F3460",
            "text_on_dark": "FFFFFF",
            "text_on_light": "2D3436",
            "muted": "636E72",
            "success": "00B894",
            "warning": "FDCB6E",
            "header_font": "Georgia",
            "body_font": "Calibri",
        },
        "warm_terracotta": {
            "name": "Warm Terracotta",
            "bg_dark": "E2725B",
            "bg_light": "FFF8F0",
            "primary": "CC5500",
            "accent": "FFDAB9",
            "text_on_dark": "FFFFFF",
            "text_on_light": "2D3436",
            "muted": "636E72",
            "success": "00B894",
            "warning": "FDCB6E",
            "header_font": "Georgia",
            "body_font": "Calibri",
        },
    }
    
    # ========================================
    # TYPOGRAPHY CONSTANTS
    # ========================================
    
    TITLE_SIZE = Pt(36)
    SUBTITLE_SIZE = Pt(20)
    HEADING_SIZE = Pt(28)
    BODY_SIZE = Pt(16)
    CAPTION_SIZE = Pt(11)
    SPEAKER_NOTE_SIZE = Pt(12)
    
    # ========================================
    # LAYOUT ENGINE
    # ========================================
    
    def select_layout(self, slide: Dict, prev_layout: str = None,
                      prev_prev_layout: str = None) -> str:
        """
        Rule-based layout selection (never 3 consecutive same layouts)
        
        Layout options:
        - title: Title slide with dark background
        - two_column: Text left (60%), image right (40%)
        - bullets_only: Full-width bullet points
        - data_callout: Large metric/number display
        - image_full: Full-width image with caption
        - diagram: Technical diagram with Mermaid.js rendering
        - chart: Data visualization slide
        """
        bullets = slide.get("bullets", [])
        has_image = bool(slide.get("image_prompt"))
        has_data = slide.get("data") is not None
        has_chart = slide.get("chart_data") is not None
        has_mermaid = slide.get("mermaid_code") is not None
        word_count = sum(len(b.split()) for b in bullets)
        
        # Hard rules
        if slide.get("slide_number") == 1:
            return "title"
        if has_mermaid:
            return "diagram"
        if has_chart:
            return "chart"
        if has_data:
            candidate = "data_callout"
        elif bullets == [] and has_image:
            candidate = "image_full"
        elif has_image and word_count < 40:
            candidate = "two_column"
        elif len(bullets) > 5:
            candidate = "bullets_only"
        elif has_image:
            candidate = "two_column"
        else:
            candidate = "bullets_only"
        
        # Anti-repetition: no 3 in a row
        if candidate == prev_layout == prev_prev_layout:
            alternatives = ["two_column", "data_callout", "bullets_only", "chart"]
            alternatives = [a for a in alternatives if a != candidate]
            candidate = alternatives[0] if alternatives else "bullets_only"
        
        return candidate
    
    def analyze_and_adjust(self, outline: Dict) -> Dict:
        """
        Pre-assembly density and layout analysis.
        
        Rules:
        - Split slides with >6 bullets
        - Enforce layout variety (no 3 consecutive same)
        - Renumber slides after adjustments
        """
        # Validate outline structure
        if not isinstance(outline, dict):
            raise ValueError(f"Outline must be a dict, got {type(outline)}")
        
        slides = outline.get("slides")
        if slides is None:
            # Log the full outline for debugging
            import logging
            logger = logging.getLogger("app.design_engine")
            logger.error(f"Outline missing 'slides' key. Keys found: {list(outline.keys())}")
            logger.error(f"Full outline: {outline}")
            raise ValueError(f"Outline missing 'slides' key. Available keys: {list(outline.keys())}")
        
        if not isinstance(slides, list):
            raise ValueError(f"'slides' must be a list, got {type(slides)}")
        adjusted = []
        prev = None
        prev_prev = None
        
        for slide in slides:
            # Density check: split slides with >6 bullets
            bullets = slide.get("bullets", [])
            if not bullets:
                # Initialize bullets if missing to prevent frontend crashes
                bullets = []
                slide["bullets"] = bullets
            
            if len(bullets) > 6:
                # Split into two slides
                mid = len(bullets) // 2
                slide1 = {**slide, "bullets": bullets[:mid],
                          "title": slide["title"] + " (1/2)"}
                slide2 = {**slide, "bullets": bullets[mid:],
                          "title": slide["title"] + " (2/2)",
                          "slide_number": slide["slide_number"] + 0.5}
                adjusted.extend([slide1, slide2])
            else:
                adjusted.append(slide)
            
            # Update layout based on rules
            slide["layout"] = self.select_layout(slide, prev, prev_prev)
            prev_prev = prev
            prev = slide["layout"]
        
        # Renumber
        for i, slide in enumerate(adjusted):
            slide["slide_number"] = i + 1
        
        outline["slides"] = adjusted
        return outline
    
    def get_theme(self, theme_name: str) -> Dict:
        """Get theme by name, fallback to ocean_gradient"""
        return self.THEMES.get(theme_name, self.THEMES["ocean_gradient"])
    
    # ========================================
    # PPTX ASSEMBLY (python-pptx)
    # ========================================
    
    def assemble_pptx(self, outline: Dict, content: List,
                      images: Dict, theme: Dict) -> bytes:
        """
        Build professional PPTX using python-pptx.
        
        Args:
            outline: Slide outline with layouts
            content: Refined slide content
            images: Dict of slide_index -> image_bytes
            theme: Theme dict from get_theme()
            
        Returns:
            PPTX file as bytes
        """
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)
        
        # Filter out empty slides (no content)
        valid_slides = []
        for i, slide_data in enumerate(outline["slides"]):
            # Check if slide has content
            has_content = (
                slide_data.get("title") or
                slide_data.get("bullets") or
                slide_data.get("supporting_data") or
                slide_data.get("chart_data") or
                slide_data.get("mermaid_code") or
                slide_data.get("subtitle_takeaway") or
                slide_data.get("speaker_notes") or
                images.get(i)  # Has an image
            )
            
            # Always keep title slides (first and last)
            is_title_slide = slide_data.get("layout") == "title"
            is_first_or_last = i == 0 or i == len(outline["slides"]) - 1
            
            if has_content or (is_title_slide and is_first_or_last):
                valid_slides.append((i, slide_data))
            else:
                print(f"⚠️ Removing empty slide {i+1}: '{slide_data.get('title', 'Untitled')}'")
        
        # Build slides
        for idx, (original_idx, slide_data) in enumerate(valid_slides):
            layout = slide_data.get("layout", "bullets_only")
            image_bytes = images.get(original_idx)
            
            if layout == "title":
                self._build_title_slide(prs, slide_data, theme, image_bytes)
            elif layout == "two_column":
                self._build_two_column_slide(prs, slide_data, theme, image_bytes)
            elif layout == "data_callout":
                self._build_data_callout_slide(prs, slide_data, theme)
            elif layout == "image_full":
                self._build_image_full_slide(prs, slide_data, theme, image_bytes)
            elif layout == "diagram":
                self._build_diagram_slide(prs, slide_data, theme)
            elif layout == "chart":
                self._build_chart_slide(prs, slide_data, theme)
            else:
                self._build_bullets_slide(prs, slide_data, theme)
            
            # Add speaker notes
            if slide_data.get("speaker_notes"):
                slide = prs.slides[-1]
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_data["speaker_notes"]
        
        # Add slide numbers to all content slides
        self._add_slide_numbers(prs, theme)
        
        # Export to bytes
        output = io.BytesIO()
        prs.save(output)
        return output.getvalue()
    
    def _build_title_slide(self, prs, data, theme, image_bytes=None):
        """Dark background, centered title, optional background image"""
        # Use blank layout (index 6) for custom design
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Dark background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor.from_string(theme["bg_dark"])

        # Background image (dimmed) if available
        if image_bytes:
            try:
                slide.shapes.add_picture(
                    io.BytesIO(image_bytes),
                    Inches(0), Inches(0),
                    prs.slide_width, prs.slide_height
                )
                # Add dark overlay
                overlay = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0), Inches(0),
                    prs.slide_width, prs.slide_height
                )
                overlay.fill.solid()
                overlay.fill.fore_color.rgb = RGBColor.from_string(theme["bg_dark"])
                overlay.fill.transparency = 0.5
            except Exception as e:
                print(f"⚠️ Title slide image failed: {e}")
        
        # Title text
        txBox = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.5), Inches(10), Inches(2)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data["title"]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(theme["text_on_dark"])
        p.font.name = theme["header_font"]
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if data.get("subtitle"):
            p2 = tf.add_paragraph()
            p2.text = data["subtitle"]
            p2.font.size = Pt(20)
            p2.font.color.rgb = RGBColor.from_string(theme["text_on_dark"])
            p2.font.name = theme["body_font"]
            p2.alignment = PP_ALIGN.CENTER
        
        # Academic metadata (Student Name, Matriculation, Course Code, Date)
        academic_meta = []
        if data.get("student_name"):
            academic_meta.append(f"Student: {data['student_name']}")
        if data.get("matriculation_number"):
            academic_meta.append(f"Matriculation: {data['matriculation_number']}")
        if data.get("course_code"):
            academic_meta.append(f"Course: {data['course_code']}")
        if data.get("date"):
            academic_meta.append(f"Date: {data['date']}")
        elif data.get("presentation_date"):
            academic_meta.append(f"Date: {data['presentation_date']}")
        
        if academic_meta:
            # Add spacing
            p_spacing = tf.add_paragraph()
            p_spacing.text = ""
            p_spacing.space_before = Pt(20)
            
            # Add academic metadata
            for meta_text in academic_meta:
                p_meta = tf.add_paragraph()
                p_meta.text = meta_text
                p_meta.font.size = Pt(14)
                p_meta.font.color.rgb = RGBColor.from_string(theme["text_on_dark"])
                p_meta.font.name = theme["body_font"]
                p_meta.alignment = PP_ALIGN.CENTER
                p_meta.space_after = Pt(4)
    
    def _build_two_column_slide(self, prs, data, theme, image_bytes=None):
        """Text left (60%), image right (40%)"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Light background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor.from_string(theme["bg_light"])
        
        # Left column: Title + bullets (60% width)
        left_x, left_w = Inches(0.8), Inches(7)
        
        # Title
        title_box = slide.shapes.add_textbox(left_x, Inches(0.5), left_w, Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data["title"]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(theme["primary"])
        p.font.name = theme["header_font"]
        
        # Bullets
        bullet_box = slide.shapes.add_textbox(left_x, Inches(1.8), left_w, Inches(5))
        tf = bullet_box.text_frame
        tf.word_wrap = True
        for j, bullet in enumerate(data.get("bullets", [])):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor.from_string(theme["text_on_light"])
            p.font.name = theme["body_font"]
            p.space_after = Pt(8)
        
        # Right column: Image (40% width)
        if image_bytes:
            try:
                img_x = Inches(8.2)
                img_w = Inches(4.8)
                img_h = Inches(5.5)
                slide.shapes.add_picture(
                    io.BytesIO(image_bytes),
                    img_x, Inches(1), img_w, img_h
                )
            except Exception as e:
                print(f"⚠️ Two-column slide image failed: {e}")
    
    def _build_data_callout_slide(self, prs, data, theme):
        """Large metric/number display"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Light background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor.from_string(theme["bg_light"])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.5), Inches(11.7), Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data["title"]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(theme["primary"])
        p.font.name = theme["header_font"]
        
        # Data value (large)
        if data.get("data"):
            value = data["data"].get("value", "")
            label = data["data"].get("label", "")
            
            value_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(2), Inches(11.7), Inches(3)
            )
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = value
            p.font.size = Pt(72)
            p.font.bold = True
            p.font.color.rgb = RGBColor.from_string(theme["primary"])
            p.alignment = PP_ALIGN.CENTER
            
            if label:
                label_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(5.5), Inches(11.7), Inches(1)
                )
                tf = label_box.text_frame
                p = tf.paragraphs[0]
                p.text = label
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor.from_string(theme["muted"])
                p.alignment = PP_ALIGN.CENTER
        
        # Bullets (if any)
        if data.get("bullets"):
            bullet_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(5), Inches(11.7), Inches(2)
            )
            tf = bullet_box.text_frame
            tf.word_wrap = True
            for j, bullet in enumerate(data["bullets"]):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor.from_string(theme["text_on_light"])
                p.font.name = theme["body_font"]
    
    def _build_image_full_slide(self, prs, data, theme, image_bytes=None):
        """Full-width image with caption"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Light background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor.from_string(theme["bg_light"])
        
        # Title
        if data.get("title"):
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.8)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = data["title"]
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor.from_string(theme["primary"])
            p.font.name = theme["header_font"]
        
        # Full-width image
        if image_bytes:
            slide.shapes.add_picture(
                io.BytesIO(image_bytes),
                Inches(0.5), Inches(1.2),
                Inches(12.333), Inches(5.5)
            )
        
        # Caption/bullets at bottom
        if data.get("bullets"):
            caption_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.6)
            )
            tf = caption_box.text_frame
            tf.word_wrap = True
            for j, bullet in enumerate(data["bullets"]):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor.from_string(theme["text_on_light"])
                p.font.name = theme["body_font"]
    
    def _build_bullets_slide(self, prs, data, theme):
        """Full-width bullet points"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Light background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor.from_string(theme["bg_light"])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.5), Inches(11.7), Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data["title"]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(theme["primary"])
        p.font.name = theme["header_font"]
        
        # Bullets
        bullet_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8), Inches(11.7), Inches(5)
        )
        tf = bullet_box.text_frame
        tf.word_wrap = True
        for j, bullet in enumerate(data.get("bullets", [])):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor.from_string(theme["text_on_light"])
            p.font.name = theme["body_font"]
            p.space_after = Pt(8)
    
    def _build_diagram_slide(self, prs, data, theme):
        """Technical diagram slide with Mermaid.js code display"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Light background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor.from_string(theme["bg_light"])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.5), Inches(11.7), Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data["title"]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(theme["primary"])
        p.font.name = theme["header_font"]
        
        # Mermaid diagram code display
        mermaid_code = data.get("mermaid_code", "")
        if mermaid_code:
            # Display diagram as code block
            diagram_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5)
            )
            tf = diagram_box.text_frame
            tf.word_wrap = True
            
            # Add code header
            p = tf.paragraphs[0]
            p.text = "[Technical Diagram - Mermaid.js]"
            p.font.size = Pt(12)
            p.font.italic = True
            p.font.color.rgb = RGBColor.from_string(theme["muted"])
            
            # Add the diagram code
            p = tf.add_paragraph()
            p.text = mermaid_code[:500] + "..." if len(mermaid_code) > 500 else mermaid_code
            p.font.size = Pt(10)
            p.font.name = "Courier New"
            p.font.color.rgb = RGBColor.from_string(theme["text_on_light"])
            p.space_before = Pt(12)
        
        # Bullets/context below diagram
        if data.get("bullets"):
            bullet_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(6.5), Inches(11.7), Inches(1)
            )
            tf = bullet_box.text_frame
            tf.word_wrap = True
            for j, bullet in enumerate(data.get("bullets", [])[:2]):  # Max 2 bullets
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor.from_string(theme["text_on_light"])
                p.font.name = theme["body_font"]
    
    def _build_chart_slide(self, prs, data, theme):
        """Chart/data visualization slide with actual chart rendering"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Light background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor.from_string(theme["bg_light"])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data["title"]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(theme["primary"])
        p.font.name = theme["header_font"]
        
        # Chart data
        chart_data = data.get("chart_data", {})
        if chart_data and chart_data.get("type") and chart_data.get("labels") and chart_data.get("values"):
            try:
                # Prepare chart data
                labels = chart_data.get("labels", [])[:8]  # Max 8 categories
                values = chart_data.get("values", [])[:8]
                chart_type = chart_data.get("type", "bar").lower()
                
                if labels and values and len(labels) == len(values):
                    # Create chart data object
                    chart_data_obj = CategoryChartData()
                    chart_data_obj.categories = labels
                    chart_data_obj.add_series('Data', values)
                    
                    # Determine chart type
                    if chart_type == "pie":
                        xl_chart_type = XL_CHART_TYPE.PIE
                    elif chart_type == "line":
                        xl_chart_type = XL_CHART_TYPE.LINE_MARKERS
                    elif chart_type == "column":
                        xl_chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
                    else:  # Default to bar
                        xl_chart_type = XL_CHART_TYPE.BAR_CLUSTERED
                    
                    # Add chart to slide
                    x, y, cx, cy = Inches(0.8), Inches(1.3), Inches(11.7), Inches(4.5)
                    chart = slide.shapes.add_chart(
                        xl_chart_type, x, y, cx, cy, chart_data_obj
                    ).chart
                    
                    # Style the chart
                    chart.has_legend = True
                    chart.legend.include_in_layout = False
                    
                    # Style the title
                    if chart.has_title:
                        chart.chart_title.text_frame.text = chart_data.get("title", "")
                else:
                    # Fallback: Display data as text if chart creation fails
                    self._display_chart_as_text(slide, chart_data, theme)
            except Exception as e:
                print(f"⚠️ Chart rendering failed: {e}. Falling back to text display.")
                self._display_chart_as_text(slide, chart_data, theme)
        
        # Context bullets below chart
        if data.get("bullets"):
            bullet_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(6), Inches(11.7), Inches(1.3)
            )
            tf = bullet_box.text_frame
            tf.word_wrap = True
            for j, bullet in enumerate(data.get("bullets", [])[:2]):  # Max 2 bullets
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor.from_string(theme["text_on_light"])
                p.font.name = theme["body_font"]
                p.space_after = Pt(4)
    
    def _display_chart_as_text(self, slide, chart_data, theme):
        """Fallback: Display chart data as formatted text"""
        chart_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8), Inches(11.7), Inches(4)
        )
        tf = chart_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"Data: {chart_data.get('type', 'Unknown').upper()} Chart"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(theme["primary"])
        
        labels = chart_data.get("labels", [])
        values = chart_data.get("values", [])
        
        if labels and values:
            for label, value in zip(labels[:8], values[:8]):
                p = tf.add_paragraph()
                p.text = f"  • {label}: {value}"
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor.from_string(theme["text_on_light"])
    
    def _add_slide_numbers(self, prs, theme):
        """Add slide numbers to content slides"""
        # Skip title slide (index 0)
        # Use range-based access to avoid prs.slides[1:] slicing bug in some pptx versions
        for i in range(1, len(prs.slides)):
            slide = prs.slides[i]
            # Add small textbox at bottom right
            num_box = slide.shapes.add_textbox(
                Inches(12), Inches(7), Inches(1), Inches(0.4)
            )
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor.from_string(theme["muted"])
            p.alignment = PP_ALIGN.RIGHT


# Singleton instance
design_engine = DesignEngine()
